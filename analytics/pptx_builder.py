"""GA4月次レポート PowerPoint生成（Shopifyレポートと同デザイン）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- カラー（generate_slides.pyと統一） ----
C_NAVY   = RGBColor(0x1A, 0x37, 0x6C)
C_BLUE   = RGBColor(0x2E, 0x86, 0xAB)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_LGRAY  = RGBColor(0xE8, 0xE8, 0xE8)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
C_TEAL   = RGBColor(0x16, 0xA0, 0x85)
C_COMMENT       = RGBColor(0xFF, 0xF8, 0xE1)
C_COMMENT_BORDER = RGBColor(0xF5, 0xA6, 0x23)

FONT    = "Meiryo"
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

CHANNEL_COLORS = {
    "Organic Search":  C_GREEN,
    "Direct":          C_NAVY,
    "Organic Social":  C_BLUE,
    "Referral":        C_ORANGE,
    "Paid Search":     C_RED,
    "Display":         C_PURPLE,
    "Email":           C_TEAL,
    "(Other)":         C_GRAY,
    "Unassigned":      C_LGRAY,
}
DEVICE_COLORS = [C_NAVY, C_BLUE, C_GRAY]

# ---- ユーティリティ ----
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=12, bold=False,
        color=C_BLACK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header(slide, title):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), C_NAVY)
    txt(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
        size=22, bold=True, color=C_WHITE)

def comment_box(slide, text, top=Inches(6.55)):
    rect(slide, Inches(0.3), top, Inches(12.73), Inches(0.75),
         C_COMMENT, C_COMMENT_BORDER)
    txt(slide, f"ポイント: {text}",
        Inches(0.45), top + Inches(0.08), Inches(12.4), Inches(0.6),
        size=10.5, color=RGBColor(0x7D, 0x4E, 0x00))

def table(slide, headers, rows, l, t, w, h, hbg=C_NAVY, alt=C_LIGHT, font_size=10):
    if not rows:
        return None
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers), l, t, w, h).table
    cw = w // len(headers)
    for i in range(len(headers)):
        tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = hbg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT; r.font.bold = True
        r.font.size = Pt(font_size + 1); r.font.color.rgb = C_WHITE
    for ri, row in enumerate(rows):
        rbg = alt if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = rbg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            r = p.runs[0]
            r.font.name = FONT; r.font.size = Pt(font_size)
            r.font.color.rgb = C_BLACK
    return tbl

# ---- フォーマッター ----
def fmt_n(v) -> str:
    return f"{int(float(v)):,}"

def fmt_pct(v) -> str:
    return f"{float(v) * 100:.1f}%"

def fmt_dur(v) -> str:
    s = int(float(v))
    return f"{s // 60}:{s % 60:02d}"

def mom_arrow(cur, prv) -> str:
    if float(prv) == 0:
        return "—"
    c = (float(cur) - float(prv)) / float(prv) * 100
    return f"{'▲' if c >= 0 else '▼'}{abs(c):.1f}%"

def mom_color(cur, prv) -> RGBColor:
    if float(prv) == 0:
        return C_GRAY
    return C_GREEN if float(cur) >= float(prv) else C_RED

# ---- スライド生成 ----

def slide_cover(prs, year: int, month: int, site_count: int):
    s = blank(prs)
    bg(s, C_WHITE)
    rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), C_NAVY)
    rect(s, Inches(0.5), Inches(2.2), Inches(0.08), Inches(2.8), C_NAVY)
    txt(s, "GA4 月次アクセスレポート",
        Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
        size=32, bold=True, color=C_NAVY)
    txt(s, f"{year}年{month}月",
        Inches(0.8), Inches(3.2), Inches(11), Inches(0.9),
        size=26, color=C_BLACK)
    txt(s, f"対象サイト数: {site_count}サイト  |  Powered by Google Analytics 4",
        Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
        size=13, color=C_GRAY)


def slide_all_summary(prs, sites_data: list[dict], year: int, month: int):
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f"全サイトサマリー — {year}年{month}月")

    rows = []
    for sd in sites_data:
        if sd.get("error"):
            rows.append([sd["name"], "取得エラー", "—", "—", "—", "—"])
            continue
        t = sd["data"]["totals"]
        p = sd["data"]["prev_totals"]
        cur = float(t.get("sessions", 0))
        prv = float(p.get("sessions", 0))
        rows.append([
            sd["name"],
            fmt_n(cur),
            mom_arrow(cur, prv),
            fmt_n(t.get("totalUsers", 0)),
            fmt_n(t.get("screenPageViews", 0)),
            fmt_pct(t.get("engagementRate", 0)),
        ])

    table(s,
          ["サイト名", "セッション", "前月比", "ユーザー", "PV", "エンゲージ率"],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5),
          font_size=10)
    comment_box(s, f"全{len(sites_data)}サイトの{year}年{month}月実績。▲=増加、▼=減少。")


def slide_site_kpi(prs, site_name: str, data: dict, year: int, month: int):
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f"{site_name} — {year}年{month}月 KPIサマリー")

    t = data["totals"]
    p = data["prev_totals"]

    kpi_items = [
        ("セッション",       fmt_n(t["sessions"]),                      C_NAVY),
        ("ユーザー",         fmt_n(t["totalUsers"]),                    C_BLUE),
        ("新規ユーザー",     fmt_n(t["newUsers"]),                      C_GREEN),
        ("ページビュー",     fmt_n(t["screenPageViews"]),               C_NAVY),
        ("平均滞在時間",     fmt_dur(t["averageSessionDuration"]),       C_ORANGE),
        ("エンゲージメント率", fmt_pct(t["engagementRate"]),             C_BLUE),
    ]
    cw = Inches(1.95); ch = Inches(1.55); gap = Inches(0.22); t0 = Inches(1.15)
    for i, (lbl, val, col) in enumerate(kpi_items):
        c = i % 3; r = i // 3
        l = Inches(0.4) + c * (cw + gap)
        ty = t0 + r * (ch + Inches(0.2))
        rect(s, l, ty, cw, ch, col)
        txt(s, lbl, l + Inches(0.1), ty + Inches(0.12),
            cw - Inches(0.2), Inches(0.4), size=10, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, val, l + Inches(0.05), ty + Inches(0.5),
            cw - Inches(0.1), Inches(0.85), size=18, bold=True,
            color=C_WHITE, align=PP_ALIGN.CENTER)

    # 右側：前月比テーブル
    rx = Inches(7.1)
    rect(s, rx, Inches(1.15), Inches(5.9), Inches(5.2), C_WHITE)
    txt(s, "前月比較", rx + Inches(0.2), Inches(1.25),
        Inches(5.5), Inches(0.4), size=12, bold=True, color=C_NAVY)

    mom_items = [
        ("セッション",   t["sessions"],       p["sessions"]),
        ("ユーザー",     t["totalUsers"],     p["totalUsers"]),
        ("新規ユーザー", t["newUsers"],       p["newUsers"]),
        ("ページビュー", t["screenPageViews"], p["screenPageViews"]),
    ]
    rh = Inches(0.95)
    for ri, (lbl, cur, prv) in enumerate(mom_items):
        ty = Inches(1.75) + ri * rh
        rbg = C_LIGHT if ri % 2 == 0 else C_WHITE
        rect(s, rx + Inches(0.1), ty, Inches(5.7), rh, rbg)
        txt(s, lbl, rx + Inches(0.25), ty + Inches(0.25),
            Inches(2.5), rh, size=11, color=C_BLACK)
        txt(s, fmt_n(cur), rx + Inches(2.5), ty + Inches(0.25),
            Inches(1.5), rh, size=11, bold=True, color=C_BLACK, align=PP_ALIGN.RIGHT)
        arrow = mom_arrow(cur, prv)
        col = mom_color(cur, prv)
        txt(s, arrow, rx + Inches(4.0), ty + Inches(0.25),
            Inches(1.5), rh, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)

    sess_arrow = mom_arrow(t["sessions"], p["sessions"])
    comment_box(s, f"セッション: {fmt_n(t['sessions'])}（前月比 {sess_arrow}）  エンゲージメント率: {fmt_pct(t['engagementRate'])}")


def slide_channels_devices(prs, site_name: str, data: dict, year: int, month: int):
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f"流入チャネル・デバイス — {site_name}  {year}年{month}月")

    channels = data["channels"]
    devices = data["devices"]
    countries = data["countries"]
    total_sessions = max(int(float(data["totals"].get("sessions", 1))), 1)

    # --- 左パネル: チャネル横棒グラフ ---
    lp_l = Inches(0.3); lp_w = Inches(6.5)
    rect(s, lp_l, Inches(1.1), lp_w, Inches(6.1), C_WHITE)
    txt(s, "流入チャネル別セッション", lp_l + Inches(0.2), Inches(1.2),
        lp_w - Inches(0.4), Inches(0.4), size=12, bold=True, color=C_NAVY)

    max_ch = max((int(float(r.get("sessions", 0))) for r in channels), default=1)
    bar_area_w = Inches(3.2)
    bar_l = lp_l + Inches(2.2)

    for i, ch in enumerate(channels[:7]):
        ch_name = ch.get("sessionDefaultChannelGroup", "Other")
        ch_sess = int(float(ch.get("sessions", 0)))
        pct = ch_sess / total_sessions * 100
        col = CHANNEL_COLORS.get(ch_name, C_GRAY)
        bt = Inches(1.75) + i * Inches(0.72)
        bw = max(bar_area_w * ch_sess / max_ch, Inches(0.05))
        rect(s, bar_l, bt, bw, Inches(0.52), col)
        txt(s, ch_name, lp_l + Inches(0.1), bt + Inches(0.07),
            Inches(2.0), Inches(0.5), size=9.5, color=C_BLACK)
        txt(s, f"{ch_sess:,}  ({pct:.1f}%)",
            bar_l + bw + Inches(0.08), bt + Inches(0.07),
            Inches(1.5), Inches(0.5), size=9.5, bold=True, color=col)

    # --- 右パネル ---
    rx = Inches(7.1)

    # デバイス内訳
    rect(s, rx, Inches(1.1), Inches(6.0), Inches(2.85), C_WHITE)
    txt(s, "デバイス内訳", rx + Inches(0.2), Inches(1.2),
        Inches(5.5), Inches(0.4), size=12, bold=True, color=C_NAVY)
    dev_max = max((int(float(r.get("sessions", 0))) for r in devices), default=1)
    for i, dev in enumerate(devices):
        dev_name = dev.get("deviceCategory", "other")
        dev_sess = int(float(dev.get("sessions", 0)))
        dev_pct = dev_sess / total_sessions * 100
        col = DEVICE_COLORS[i % len(DEVICE_COLORS)]
        bt = Inches(1.73) + i * Inches(0.7)
        bw = max(Inches(3.0) * dev_sess / dev_max, Inches(0.05))
        bl = rx + Inches(1.5)
        rect(s, bl, bt, bw, Inches(0.5), col)
        txt(s, dev_name, rx + Inches(0.1), bt + Inches(0.07),
            Inches(1.3), Inches(0.5), size=10, color=C_BLACK)
        txt(s, f"{dev_sess:,}  ({dev_pct:.1f}%)",
            bl + bw + Inches(0.1), bt + Inches(0.07),
            Inches(1.8), Inches(0.5), size=10, bold=True, color=col)

    # 国別TOP5
    rect(s, rx, Inches(4.1), Inches(6.0), Inches(2.85), C_WHITE)
    txt(s, "アクセス国別 TOP5", rx + Inches(0.2), Inches(4.2),
        Inches(5.5), Inches(0.4), size=12, bold=True, color=C_NAVY)
    ctry_colors = [C_NAVY, C_BLUE, C_ORANGE, C_GREEN, C_RED]
    ctry_max = max(int(float(countries[0].get("sessions", 1))), 1) if countries else 1
    for i, ctry in enumerate(countries[:5]):
        ctry_name = ctry.get("country", "Unknown")
        ctry_sess = int(float(ctry.get("sessions", 0)))
        ctry_pct = ctry_sess / total_sessions * 100
        col = ctry_colors[i % len(ctry_colors)]
        bt = Inches(4.73) + i * Inches(0.43)
        bw = max(Inches(3.0) * ctry_sess / ctry_max, Inches(0.05))
        bl = rx + Inches(1.8)
        rect(s, bl, bt, bw, Inches(0.35), col)
        txt(s, ctry_name, rx + Inches(0.1), bt + Inches(0.02),
            Inches(1.6), Inches(0.38), size=9, color=C_BLACK)
        txt(s, f"{ctry_sess:,} ({ctry_pct:.1f}%)",
            bl + bw + Inches(0.08), bt + Inches(0.02),
            Inches(1.6), Inches(0.38), size=9, bold=True, color=col)

    top_ch = channels[0] if channels else {}
    top_name = top_ch.get("sessionDefaultChannelGroup", "—")
    top_pct = int(float(top_ch.get("sessions", 0))) / total_sessions * 100 if channels else 0
    comment_box(s, f"主要流入: {top_name}（{top_pct:.1f}%）  デバイスやチャネル構成を前月と比較してください。")


def build_report(sites_data: list[dict], year: int, month: int, output_path: str):
    prs = new_prs()
    slide_cover(prs, year, month, len(sites_data))
    slide_all_summary(prs, sites_data, year, month)
    for sd in sites_data:
        if sd.get("error"):
            continue
        slide_site_kpi(prs, sd["name"], sd["data"], year, month)
        slide_channels_devices(prs, sd["name"], sd["data"], year, month)
    prs.save(output_path)
    print(f"OK: {output_path}")
