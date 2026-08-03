# -*- coding: utf-8 -*-
"""
GA4 詳細サイトレポート生成
使い方: python detail_report.py [--site サイト名] [--month YYYY-MM]
"""

import argparse, calendar, io, os, pickle, statistics, sys
from datetime import date
from pathlib import Path
import yaml

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker
    import matplotlib.font_manager as _fm
    import numpy as np
    HAS_MPL = True
    _available = [f.name for f in _fm.fontManager.ttflist]
    for _candidate in ['Noto Sans JP', 'Hiragino Kaku Gothic Pro', 'Meiryo', 'Yu Gothic']:
        if _candidate in _available:
            plt.rcParams['font.family'] = _candidate
            break
    plt.rcParams['axes.unicode_minus'] = False
except ImportError:
    HAS_MPL = False
    print('警告: matplotlib 未インストール。pip install matplotlib numpy で追加してください。')
    sys.exit(1)

sys.path.insert(0, str(Path(__file__).parent))
from ga_client import run_report, fetch_funnel_data

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx_builder import (
    blank, bg, rect, txt, header, table, new_prs, comment_box,
    fmt_n, fmt_pct, fmt_dur, mom_arrow, mom_color,
    C_NAVY, C_BLUE, C_LIGHT, C_WHITE, C_BLACK, C_GRAY,
    C_GREEN, C_RED, C_ORANGE, C_PURPLE, C_TEAL, C_LGRAY,
    SLIDE_W, SLIDE_H, FONT, C_COMMENT, C_COMMENT_BORDER,
)

BASE_DIR   = Path(__file__).parent
CONFIG_YML = BASE_DIR / 'config.yaml'

CHART_COLORS = ['#2272B5', '#F07855', '#3BAD9B', '#8B7DB8', '#E8C96A',
                '#5096C8', '#E88A78', '#4BAD8A', '#A09BC0', '#D4EAF7']

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ユーティリティ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def get_date_range(year, month):
    last = calendar.monthrange(year, month)[1]
    return (date(year, month, 1).strftime('%Y-%m-%d'),
            date(year, month, last).strftime('%Y-%m-%d'))

def prev_month(year, month):
    return (year - 1, 12) if month == 1 else (year, month - 1)

def safe_run(property_id, start, end, metrics, dims=None, order=None, limit=10):
    try:
        return run_report(property_id, start, end, metrics, dims, order, limit)
    except Exception as ex:
        print(f'  [skip] dims={dims} metrics={metrics}: {ex}')
        return []

LOW_BENCHMARK_RATIO = 0.7  # 他サイト中央値の70%未満を「目立って低い」と判定

def fetch_benchmark_metrics(sites, year, month):
    """全サイト共通の月次指標（CVR・エンゲージメント率）を取得し、サイト間比較に使う"""
    start, end = get_date_range(year, month)
    benchmark = {}
    for s in sites:
        name = s['name']
        pid = str(s['property_id'])
        rows = safe_run(pid, start, end, ['sessions', 'engagedSessions'])
        t = rows[0] if rows else {}
        sessions = float(t.get('sessions', 0) or 0)
        engaged  = float(t.get('engagedSessions', 0) or 0)
        engagement_rate = engaged / sessions * 100 if sessions > 0 else None

        cvr = None
        if s.get('has_ecommerce') and sessions > 0:
            funnel = fetch_funnel_data(pid, start, end)
            cvr = funnel.get('purchase', 0) / sessions * 100

        benchmark[name] = {'engagement_rate': engagement_rate, 'cvr': cvr}
    return benchmark

def benchmark_lowlights(site_name, benchmark):
    """他サイトの中央値と比較して目立って低い指標のコメント文を返す"""
    site = benchmark.get(site_name, {})
    labels = [('cvr', 'CVR（購入完了率）'), ('engagement_rate', 'エンゲージメント率')]
    comments = []
    for key, label in labels:
        value = site.get(key)
        if value is None:
            continue
        others = [v[key] for k, v in benchmark.items() if k != site_name and v.get(key) is not None]
        if len(others) < 3:
            continue
        med = statistics.median(others)
        if med > 0 and value < med * LOW_BENCHMARK_RATIO:
            gap = (1 - value / med) * 100
            comments.append(f'・{label}: {value:.2f}%（他サイト中央値 {med:.2f}% より{gap:.0f}%低い）')
    return comments

def fetch_search_console(site_url, start_date, end_date, limit=10):
    """Search Console API からクエリデータを取得（OAuth ユーザー認証）"""
    try:
        from google_auth_oauthlib.flow import InstalledAppFlow
        from google.auth.transport.requests import Request as GoogleRequest
        from googleapiclient.discovery import build
    except ImportError:
        print('  [skip] google-api-python-client または google-auth-oauthlib が未インストール')
        return []

    oauth_path  = BASE_DIR / 'oauth_client.json'
    token_path  = BASE_DIR / 'gsc_token.pickle'
    SCOPES_GSC  = ['https://www.googleapis.com/auth/webmasters.readonly']

    if not oauth_path.exists():
        print(f'  [skip] oauth_client.json が見つかりません: {oauth_path}')
        return []

    creds = None
    if token_path.exists():
        with open(token_path, 'rb') as f:
            creds = pickle.load(f)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(GoogleRequest())
            except Exception:
                creds = None
        if not creds:
            flow = InstalledAppFlow.from_client_secrets_file(str(oauth_path), SCOPES_GSC)
            creds = flow.run_local_server(port=0)
        with open(token_path, 'wb') as f:
            pickle.dump(creds, f)

    try:
        service  = build('webmasters', 'v3', credentials=creds)
        response = service.searchanalytics().query(
            siteUrl=site_url,
            body={
                'startDate':  start_date,
                'endDate':    end_date,
                'dimensions': ['query'],
                'rowLimit':   limit,
                'orderBy': [{'fieldName': 'clicks', 'sortOrder': 'DESCENDING'}],
            }
        ).execute()
        return response.get('rows', [])
    except Exception as ex:
        print(f'  [skip] Search Console API エラー: {ex}')
        return []


def chart_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    plt.close(fig)
    return buf

def frev(v):
    v = float(v or 0)
    return f'¥{int(v):,}' if v > 0 else '—'

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# データ取得
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def fetch_detail_data(property_id, year, month):
    start, end   = get_date_range(year, month)
    py, pm       = prev_month(year, month)
    pstart, pend = get_date_range(py, pm)
    pid = property_id

    yoy_start, yoy_end = get_date_range(year - 1, month)

    print('  基本集計...')
    t_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    totals = t_rows[0] if t_rows else {}

    p_rows = safe_run(pid, pstart, pend,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    prev_totals = p_rows[0] if p_rows else {}

    y_rows = safe_run(pid, yoy_start, yoy_end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    yoy_totals = y_rows[0] if y_rows else {}

    print('  デバイス別日次...')
    device_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'deviceCategory'], limit=200)
    device_daily.sort(key=lambda r: r.get('date', ''))

    print('  ユーザー属性...')
    gender_rows = safe_run(pid, start, end,
        ['sessions'], ['userGender'], order='sessions', limit=5)
    age_rows = safe_run(pid, start, end,
        ['sessions'], ['userAgeBracket'], order='sessions', limit=10)
    retention_rows = safe_run(pid, start, end,
        ['sessions'], ['newVsReturning'], order='sessions', limit=5)

    print('  ページタイトル...')
    page_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions'],
        ['pageTitle'], order='sessions', limit=10)

    print('  参照元/メディア...')
    source_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'],
        ['sessionSourceMedium'], order='sessions', limit=10)

    print('  キャンペーン...')
    campaign_rows = safe_run(pid, start, end,
        ['sessions', 'totalRevenue'],
        ['sessionCampaignName', 'sessionSourceMedium'], order='sessions', limit=10)

    print('  日次収益...')
    revenue_daily = safe_run(pid, start, end,
        ['totalRevenue', 'sessions'], ['date'], limit=35)
    revenue_daily.sort(key=lambda r: r.get('date', ''))

    print('  アイテム別収益...')
    item_rows = safe_run(pid, start, end,
        ['itemsPurchased', 'itemRevenue'],
        ['itemName'], order='itemRevenue', limit=10)

    print('  チャネル別日次...')
    channel_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'sessionDefaultChannelGroup'], limit=500)

    print('  参照元別日次...')
    source_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'sessionSourceMedium'], limit=1000)

    print('  曜日別...')
    dow_rows = safe_run(pid, start, end,
        ['sessions', 'totalRevenue'], ['dayOfWeek'], limit=7)
    dow_rows.sort(key=lambda r: int(r.get('dayOfWeek', 0)))

    print('  直近3ヶ月推移...')
    monthly_trend = []
    cy, cm = year, month
    for _ in range(3):
        ms, me = get_date_range(cy, cm)
        rows = safe_run(pid, ms, me, ['sessions', 'totalRevenue'], limit=1)
        r = rows[0] if rows else {}
        monthly_trend.insert(0, {
            'year': cy, 'month': cm,
            'sessions': float(r.get('sessions', 0)),
            'revenue':  float(r.get('totalRevenue', 0)),
        })
        cy, cm = prev_month(cy, cm)

    print('  購買ファネル...')
    try:
        funnel = fetch_funnel_data(pid, start, end)
    except Exception as ex:
        print(f'  [skip] 購買ファネル: {ex}')
        funnel = {'view_item': 0, 'add_to_cart': 0, 'begin_checkout': 0, 'purchase': 0}

    return {
        'totals':       totals,
        'prev_totals':  prev_totals,
        'yoy_totals':   yoy_totals,
        'device_daily': device_daily,
        'gender':       gender_rows,
        'age':          age_rows,
        'retention':    retention_rows,
        'pages':        page_rows,
        'sources':      source_rows,
        'campaigns':    campaign_rows,
        'revenue_daily': revenue_daily,
        'items':        item_rows,
        'dow':          dow_rows,
        'monthly_trend': monthly_trend,
        'channel_daily': channel_daily,
        'source_daily':  source_daily,
        'funnel':        funnel,
    }

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# チャート生成
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _no_data_fig(w=9, h=3.5, msg='データなし'):
    fig, ax = plt.subplots(figsize=(w, h))
    ax.text(0.5, 0.5, msg, ha='center', va='center', transform=ax.transAxes,
            fontsize=12, color='#888888')
    ax.axis('off')
    fig.tight_layout()
    return chart_to_image(fig)

def chart_device_line(device_daily):
    if not device_daily:
        return _no_data_fig(9, 3.5)
    devices = sorted(set(r['deviceCategory'] for r in device_daily))
    dates   = sorted(set(r['date'] for r in device_daily))
    dmap    = {}
    for r in device_daily:
        dmap.setdefault(r['deviceCategory'], {})[r['date']] = float(r.get('sessions', 0))
    colors_d = {'desktop': '#2272B5', 'mobile': '#F07855', 'tablet': '#3BAD9B'}
    labels_d = {'desktop': 'デスクトップ', 'mobile': 'モバイル', 'tablet': 'タブレット'}

    fig, ax = plt.subplots(figsize=(9, 3.5))
    for dev in devices:
        vals = [dmap.get(dev, {}).get(d, 0) for d in dates]
        ax.plot(range(len(dates)), vals,
                label=labels_d.get(dev, dev),
                color=colors_d.get(dev, '#999999'),
                linewidth=2, marker='o', markersize=3)
    step = max(1, len(dates) // 10)
    ax.set_xticks(range(0, len(dates), step))
    ax.set_xticklabels([dates[i][5:] for i in range(0, len(dates), step)], fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'))
    ax.set_ylabel('セッション数', fontsize=9)
    ax.legend(loc='upper right', fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('デバイス別セッション数（日次）', fontsize=10, fontweight='bold', pad=6)
    fig.tight_layout()
    return chart_to_image(fig)

def chart_pie(rows, dim_key, label_map=None, title='', colors=None, figsize=(4.2, 4.2)):
    valid = [(r.get(dim_key, ''), float(r.get('sessions', 0))) for r in rows
             if r.get(dim_key, '') not in ('', '(not set)', 'unknown')]
    if not valid:
        fig, ax = plt.subplots(figsize=figsize)
        ax.text(0.5, 0.5, 'データなし\n(デモグラフィクス未設定)',
                ha='center', va='center', fontsize=9, color='#888888')
        ax.axis('off')
        if title:
            ax.set_title(title, fontsize=10, fontweight='bold')
        fig.tight_layout()
        return chart_to_image(fig)

    labels = [(label_map or {}).get(k, k) for k, _ in valid]
    values = [v for _, v in valid]
    if not colors:
        colors = CHART_COLORS[:len(labels)]
    fig, ax = plt.subplots(figsize=figsize)
    wedges, _, autotexts = ax.pie(
        values, autopct='%1.1f%%', colors=colors[:len(labels)],
        startangle=90, pctdistance=0.72)
    for at in autotexts:
        at.set_fontsize(8)
    ax.legend(wedges, labels, loc='upper center', bbox_to_anchor=(0.5, -0.06),
              fontsize=8, framealpha=0.9, ncol=min(3, len(labels)))
    if title:
        ax.set_title(title, fontsize=10, fontweight='bold', pad=8)
    fig.tight_layout()
    return chart_to_image(fig)

def chart_revenue_line(revenue_daily):
    if not revenue_daily:
        return _no_data_fig(11, 2.8)
    dates    = [r['date'] for r in revenue_daily]
    revenues = [float(r.get('totalRevenue', 0)) for r in revenue_daily]
    fig, ax = plt.subplots(figsize=(11, 2.8))
    ax.fill_between(range(len(dates)), revenues, alpha=0.18, color='#2272B5')
    ax.plot(range(len(dates)), revenues, color='#2272B5', linewidth=2)
    step = max(1, len(dates) // 10)
    ax.set_xticks(range(0, len(dates), step))
    ax.set_xticklabels([dates[i][5:] for i in range(0, len(dates), step)], fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'¥{int(x):,}'))
    ax.set_ylabel('収益（円）', fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('日次収益', fontsize=10, fontweight='bold', pad=6)
    fig.tight_layout()
    return chart_to_image(fig)

def chart_dow_bar(dow_rows):
    if not dow_rows:
        return _no_data_fig(5.5, 2.8)
    # GA4: 0=Sun, 1=Mon ... 6=Sat → 月-日 順に並べ替え
    order  = ['1', '2', '3', '4', '5', '6', '0']
    labels = ['月', '火', '水', '木', '金', '土', '日']
    dmap   = {r.get('dayOfWeek', '0'): float(r.get('sessions', 0)) for r in dow_rows}
    values = [dmap.get(k, 0) for k in order]
    colors = ['#2272B5'] * 5 + ['#5096C8', '#F07855']

    fig, ax = plt.subplots(figsize=(5.5, 2.8))
    bars = ax.bar(labels, values, color=colors, edgecolor='white', linewidth=0.5)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'))
    ax.set_ylabel('セッション', fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('曜日別セッション数', fontsize=10, fontweight='bold', pad=6)
    for bar, val in zip(bars, values):
        if val > 0:
            ax.text(bar.get_x() + bar.get_width() / 2, val,
                    f'{int(val):,}', ha='center', va='bottom', fontsize=7)
    fig.tight_layout()
    return chart_to_image(fig)

def chart_monthly_trend(monthly_trend):
    if not monthly_trend:
        return _no_data_fig(9, 4.0)
    labels   = [f"{r['year']}年{r['month']}月" for r in monthly_trend]
    sessions = [r['sessions'] for r in monthly_trend]
    revenues = [r['revenue']  for r in monthly_trend]
    x = np.arange(len(labels))
    w = 0.35
    fig, ax1 = plt.subplots(figsize=(9, 4.0))
    ax1.bar(x - w/2, sessions, w, label='セッション数', color='#2272B5', alpha=0.85)
    ax1.set_xticks(x); ax1.set_xticklabels(labels, fontsize=11)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f'{int(v):,}'))
    ax1.set_ylabel('セッション数', color='#2272B5', fontsize=9)
    ax1.tick_params(axis='y', labelcolor='#2272B5')
    ax2 = ax1.twinx()
    ax2.bar(x + w/2, revenues, w, label='収益', color='#F07855', alpha=0.85)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f'¥{int(v):,}'))
    ax2.set_ylabel('収益（円）', color='#C0512A', fontsize=9)
    ax2.tick_params(axis='y', labelcolor='#C0512A')
    h1, l1 = ax1.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax1.legend(h1 + h2, l1 + l2, loc='upper left', fontsize=9)
    ax1.grid(axis='y', alpha=0.2)
    ax1.spines['top'].set_visible(False)
    ax1.set_title('直近3ヶ月推移（セッション数・収益）', fontsize=10, fontweight='bold', pad=6)
    fig.tight_layout()
    return chart_to_image(fig)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スパイク分析
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

CHANNEL_JP = {
    'Organic Search':    'オーガニック検索',
    'Direct':            'ダイレクト',
    'Email':             'メール',
    'Paid Search':       '有料検索',
    'Organic Social':    'オーガニックSNS',
    'Paid Social':       '有料SNS',
    'Social':            'SNS',
    'Referral':          '参照元サイト',
    'Affiliates':        'アフィリエイト',
    'Display':           'ディスプレイ広告',
    'Organic Shopping':  'オーガニックショッピング',
    'Organic Video':     'オーガニック動画',
    'Unassigned':        '未分類',
    '(Other)':           'その他',
}

def detect_spikes(device_daily, percentile=80):
    """月次日次データから急増日を検出（上位 percentile% に入る日）"""
    daily_totals = {}
    for r in device_daily:
        d = r.get('date', '')
        daily_totals[d] = daily_totals.get(d, 0) + float(r.get('sessions', 0))
    if not daily_totals:
        return [], 0.0
    values = list(daily_totals.values())
    baseline       = float(np.median(values))
    threshold_val  = float(np.percentile(values, percentile))
    spikes = [
        {'date': d, 'sessions': v, 'ratio': v / baseline}
        for d, v in sorted(daily_totals.items())
        if v >= threshold_val
    ]
    spikes.sort(key=lambda x: x['ratio'], reverse=True)
    return spikes[:5], baseline


def analyze_spike_causes(spike_date, channel_daily):
    """チャネル別日次データからスパイク日の原因チャネルを特定"""
    channel_by_date = {}
    channels_all = set()
    for r in channel_daily:
        d  = r.get('date', '')
        ch = r.get('sessionDefaultChannelGroup', 'Unassigned')
        channel_by_date.setdefault(d, {})[ch] = float(r.get('sessions', 0))
        channels_all.add(ch)
    if not channel_by_date:
        return []
    other_dates = [d for d in channel_by_date if d != spike_date]
    if not other_dates:
        return []
    avg_by_ch = {
        ch: sum(channel_by_date[d].get(ch, 0) for d in other_dates) / len(other_dates)
        for ch in channels_all
    }
    spike_channels = channel_by_date.get(spike_date, {})
    analysis = []
    for ch, spike_sess in spike_channels.items():
        avg_sess = avg_by_ch.get(ch, 0)
        delta    = spike_sess - avg_sess
        ratio    = spike_sess / avg_sess if avg_sess > 0 else spike_sess
        analysis.append({
            'channel':    ch,
            'channel_jp': CHANNEL_JP.get(ch, ch),
            'sessions':   spike_sess,
            'avg':        avg_sess,
            'delta':      delta,
            'ratio':      ratio,
        })
    analysis.sort(key=lambda x: x['delta'], reverse=True)
    return analysis[:4]


def build_spike_comment(analysis):
    """原因分析の要約コメントを生成"""
    if not analysis:
        return '特定チャネルへの集中は確認されませんでした'
    parts = []
    for a in analysis[:3]:
        if a['ratio'] >= 3.0:
            parts.append(f"{a['channel_jp']}が通常の{a['ratio']:.1f}倍")
        elif a['ratio'] >= 1.5:
            parts.append(f"{a['channel_jp']}が+{int((a['ratio'] - 1) * 100)}%増")
        elif a['delta'] > 30:
            parts.append(f"{a['channel_jp']}が増加（+{int(a['delta'])}件）")
    return ' ／ '.join(parts) if parts else '複数チャネルで均等増加'


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド生成
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _kpi_card(s, l, t, w, h, label, value, arrow, col, prev_label='前月比',
              yoy_arrow=None, yoy_label='前年同月比'):
    rect(s, l, t, w, h, col)
    if yoy_arrow is not None:
        txt(s, label,
            l + Inches(0.1), t + Inches(0.06), w - Inches(0.2), Inches(0.26),
            size=7.5, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, value,
            l + Inches(0.05), t + Inches(0.31), w - Inches(0.1), Inches(0.44),
            size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, f'{prev_label}  {arrow}',
            l + Inches(0.05), t + Inches(0.74), w - Inches(0.1), Inches(0.22),
            size=7, color=C_LGRAY, align=PP_ALIGN.CENTER)
        txt(s, f'{yoy_label}  {yoy_arrow}',
            l + Inches(0.05), t + Inches(0.96), w - Inches(0.1), Inches(0.22),
            size=7, color=C_LGRAY, align=PP_ALIGN.CENTER)
    else:
        txt(s, label,
            l + Inches(0.1), t + Inches(0.1), w - Inches(0.2), Inches(0.35),
            size=9, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, value,
            l + Inches(0.05), t + Inches(0.42), w - Inches(0.1), Inches(0.52),
            size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, f'{prev_label} {arrow}',
            l + Inches(0.05), t + Inches(0.91), w - Inches(0.1), Inches(0.28),
            size=8, color=C_LGRAY, align=PP_ALIGN.CENTER)


def slide_cover(prs, site_name, year, month):
    s = blank(prs)
    bg(s, C_WHITE)
    rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), C_NAVY)
    rect(s, Inches(0.5), Inches(1.8), Inches(0.08), Inches(3.5), C_NAVY)
    txt(s, site_name,
        Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.1),
        size=28, bold=True, color=C_NAVY)
    txt(s, 'GA4 詳細アクセス・収益レポート',
        Inches(0.8), Inches(2.9), Inches(11), Inches(0.8),
        size=20, color=C_BLUE)
    txt(s, f'{year}年{month}月',
        Inches(0.8), Inches(3.7), Inches(11), Inches(0.7),
        size=18, color=C_BLACK)
    txt(s, 'Powered by Google Analytics 4',
        Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
        size=11, color=C_GRAY)


def slide_traffic_summary(prs, site_name, data, year, month, period_label=None):
    """集客サマリ：KPI x3 + デバイス別折れ線グラフ"""
    s = blank(prs)
    bg(s, C_LIGHT)
    _pl = period_label or f'{year}年{month}月'
    header(s, f'集客サマリ — {site_name}  {_pl}')

    t = data['totals'];  p = data['prev_totals'];  y = data.get('yoy_totals', {})
    pl_kpi  = '前期比'     if period_label else '前月比'
    yoy_lbl = '前年同期比' if period_label else '前年同月比'
    kpis = [
        ('セッション数',   fmt_n(t.get('sessions', 0)),
         mom_arrow(t.get('sessions', 0),   p.get('sessions', 0)),
         mom_arrow(t.get('sessions', 0),   y.get('sessions', 0))   if y else None,
         C_NAVY),
        ('総ユーザー数',   fmt_n(t.get('totalUsers', 0)),
         mom_arrow(t.get('totalUsers', 0),  p.get('totalUsers', 0)),
         mom_arrow(t.get('totalUsers', 0),  y.get('totalUsers', 0)) if y else None,
         C_BLUE),
        ('新規ユーザー数', fmt_n(t.get('newUsers', 0)),
         mom_arrow(t.get('newUsers', 0),    p.get('newUsers', 0)),
         mom_arrow(t.get('newUsers', 0),    y.get('newUsers', 0))   if y else None,
         C_GREEN),
    ]
    cw = Inches(3.9); gap = Inches(0.25)
    for i, (lbl, val, arrow, yoy_arrow, col) in enumerate(kpis):
        _kpi_card(s, Inches(0.3) + i * (cw + gap), Inches(1.1), cw, Inches(1.22),
                  lbl, val, arrow, col, prev_label=pl_kpi,
                  yoy_arrow=yoy_arrow, yoy_label=yoy_lbl)

    img = chart_device_line(data['device_daily'])
    s.shapes.add_picture(img, Inches(0.3), Inches(2.45), Inches(12.7), Inches(4.0))


def slide_demographics(prs, site_name, data, year, month, period_label=None):
    """ユーザー属性：男女別・年齢層別・新規/リピーター 円グラフ"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'ユーザー属性 — {site_name}  {period_label or f"{year}年{month}月"}')

    gender_map = {'male': '男性', 'female': '女性'}
    age_map = {
        '18-24': '18〜24歳', '25-34': '25〜34歳', '35-44': '35〜44歳',
        '45-54': '45〜54歳', '55-64': '55〜64歳', '65+': '65歳以上',
    }
    retention_map = {'new': '新規ユーザー', 'returning': 'リピーター'}

    img_g = chart_pie(data['gender'], 'userGender', gender_map, '男女別',
                      ['#2272B5', '#F07855'])
    img_a = chart_pie(data['age'], 'userAgeBracket', age_map, '年齢層別',
                      CHART_COLORS)
    img_r = chart_pie(data.get('retention', []), 'newVsReturning', retention_map,
                      '新規 / リピーター', ['#5096C8', '#E8C96A'])

    chart_w = Inches(4.1)
    # width のみ指定してアスペクト比を維持（縦伸び防止）
    s.shapes.add_picture(img_g, Inches(0.3), Inches(1.2), width=chart_w)
    s.shapes.add_picture(img_a, Inches(4.6), Inches(1.2), width=chart_w)
    s.shapes.add_picture(img_r, Inches(8.9), Inches(1.2), width=chart_w)

    txt(s,
        '※ データなしの場合は GA4 管理画面 → データ設定 → データ収集 で「Googleシグナルのデータ収集」を有効化してください',
        Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.5),
        size=8.5, color=C_GRAY, italic=True)


def slide_page_traffic(prs, site_name, data, year, month, period_label=None):
    """ページタイトル別トラフィック TOP10"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'ページタイトル別トラフィック TOP10 — {site_name}  {period_label or f"{year}年{month}月"}')
    rows = []
    for r in data['pages']:
        rows.append([
            r.get('pageTitle', '—')[:55],
            fmt_n(r.get('sessions', 0)),
            fmt_n(r.get('totalUsers', 0)),
            fmt_n(r.get('newUsers', 0)),
            fmt_n(r.get('engagedSessions', 0)),
        ])
    if not rows:
        txt(s, 'データなし', Inches(0.3), Inches(2.5), Inches(12.7), Inches(1.0),
            size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
        return
    table(s,
          ['ページタイトル', 'セッション', '総ユーザー', '新規ユーザー', 'エンゲージメント'],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
          font_size=9)


def slide_source_medium(prs, site_name, data, year, month, period_label=None):
    """トラフィック獲得（参照元/メディア）TOP10"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'トラフィック獲得（参照元 / メディア）TOP10 — {site_name}  {period_label or f"{year}年{month}月"}')
    rows = []
    for r in data['sources']:
        rows.append([
            r.get('sessionSourceMedium', '—')[:35],
            fmt_n(r.get('sessions', 0)),
            fmt_n(r.get('totalUsers', 0)),
            fmt_n(r.get('newUsers', 0)),
            fmt_n(r.get('engagedSessions', 0)),
            frev(r.get('totalRevenue', 0)),
        ])
    if not rows:
        txt(s, 'データなし', Inches(0.3), Inches(2.5), Inches(12.7), Inches(1.0),
            size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
        return
    table(s,
          ['参照元 / メディア', 'セッション', '総ユーザー', '新規ユーザー', 'エンゲージメント', '合計収益'],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
          font_size=9)


def slide_campaign(prs, site_name, data, year, month, period_label=None):
    """トラフィック獲得（キャンペーン）TOP10"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'トラフィック獲得（キャンペーン）TOP10 — {site_name}  {period_label or f"{year}年{month}月"}')
    rows = []
    for r in data['campaigns']:
        name   = r.get('sessionCampaignName', '—') or '（未設定）'
        source = r.get('sessionSourceMedium', '—') or '（未設定）'
        rows.append([
            name[:45],
            source[:35],
            fmt_n(r.get('sessions', 0)),
            frev(r.get('totalRevenue', 0)),
        ])
    if not rows:
        txt(s, 'データなし', Inches(0.3), Inches(2.5), Inches(12.7), Inches(1.0),
            size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
        return
    table(s,
          ['キャンペーン名', '参照元/メディア', 'セッション数', '合計収益'],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
          font_size=9,
          col_widths=[Inches(4.5), Inches(3.8), Inches(2.2), Inches(2.2)])


def slide_search_query(prs, site_name, year, month, gsc_rows=None, period_label=None):
    """検索クエリ TOP10（Search Console データ）"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'検索クエリ TOP10 — {site_name}  {period_label or f"{year}年{month}月"}')

    if not gsc_rows:
        rect(s, Inches(1.8), Inches(2.3), Inches(9.7), Inches(3.0),
             C_COMMENT, C_COMMENT_BORDER)
        txt(s, 'Google Search Console との連携が必要です',
            Inches(2.0), Inches(2.55), Inches(9.3), Inches(0.65),
            size=15, bold=True, color=RGBColor(0x2A, 0x5F, 0x8A), align=PP_ALIGN.CENTER)
        txt(s,
            ('config.yaml の search_console_url と oauth_client.json を設定すると\n'
             'クエリ・クリック数・表示回数・CTR・掲載順位を自動取得できます。'),
            Inches(2.0), Inches(3.2), Inches(9.3), Inches(1.2),
            size=11, color=RGBColor(0x2A, 0x5F, 0x8A), align=PP_ALIGN.CENTER)
        return

    rows = []
    for r in gsc_rows:
        query = r.get('keys', ['—'])[0]
        rows.append([
            query[:60],
            f'{int(r.get("clicks", 0)):,}',
            f'{int(r.get("impressions", 0)):,}',
            f'{r.get("ctr", 0) * 100:.1f}%',
            f'{r.get("position", 0):.1f}',
        ])
    table(s,
          ['クエリ', 'クリック数', '表示回数', 'CTR', '平均掲載順位'],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
          font_size=10)


def slide_revenue_summary(prs, site_name, data, year, month):
    """収益サマリ：合計収益KPI + 日次グラフ + 曜日別棒グラフ"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'収益サマリ — {site_name}  {year}年{month}月')

    total_rev = float(data['totals'].get('totalRevenue', 0))
    prev_rev  = float(data['prev_totals'].get('totalRevenue', 0))
    _kpi_card(s, Inches(0.3), Inches(1.1), Inches(3.6), Inches(1.22),
              '合計収益（税込）', frev(total_rev),
              mom_arrow(total_rev, prev_rev), C_NAVY)

    # 日次収益グラフ（右寄り）
    img_rev = chart_revenue_line(data['revenue_daily'])
    s.shapes.add_picture(img_rev, Inches(4.1), Inches(1.1), Inches(9.0), Inches(2.5))

    # 曜日別棒グラフ（左下）
    img_dow = chart_dow_bar(data['dow'])
    s.shapes.add_picture(img_dow, Inches(0.3), Inches(3.8), Inches(6.0), Inches(3.1))

    # 余白に補足テキスト
    txt(s, '曜日別傾向から施策タイミングの参考にしてください。',
        Inches(6.5), Inches(3.9), Inches(6.5), Inches(0.5),
        size=10, color=C_GRAY)


def slide_item_revenue(prs, site_name, data, year, month):
    """アイテム別収益ランキング TOP10（単体スライド・後方互換用）"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'アイテム別収益ランキング TOP10 — {site_name}  {year}年{month}月')
    rows = []
    for i, r in enumerate(data['items'], 1):
        rows.append([
            str(i),
            r.get('itemName', '—')[:55],
            fmt_n(r.get('itemsPurchased', 0)),
            frev(r.get('itemRevenue', 0)),
        ])
    if not rows:
        txt(s, 'データなし（eコマース未設定または収益ゼロ）',
            Inches(0.3), Inches(2.5), Inches(12.7), Inches(1.0),
            size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
        return
    table(s,
          ['順位', 'アイテム名', '購入数', 'アイテム収益'],
          rows,
          Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.8),
          font_size=10)


def slide_revenue_combined(prs, site_name, data, year, month, period_label=None):
    """収益サマリ + アイテム別収益ランキング（1スライドに統合）"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'収益サマリ / アイテム別収益 — {site_name}  {period_label or f"{year}年{month}月"}')

    total_rev = float(data['totals'].get('totalRevenue', 0))
    prev_rev  = float(data['prev_totals'].get('totalRevenue', 0))
    yoy_rev   = float(data.get('yoy_totals', {}).get('totalRevenue', 0))
    pl_kpi  = '前期比'     if period_label else '前月比'
    yoy_lbl = '前年同期比' if period_label else '前年同月比'

    # KPI カード（左上）
    _kpi_card(s, Inches(0.3), Inches(1.1), Inches(2.5), Inches(1.22),
              '合計収益（税込）', frev(total_rev), mom_arrow(total_rev, prev_rev), C_NAVY,
              prev_label=pl_kpi,
              yoy_arrow=mom_arrow(total_rev, yoy_rev) if yoy_rev else None,
              yoy_label=yoy_lbl)

    # 日次収益グラフ（KPI の右）
    img_rev = chart_revenue_line(data['revenue_daily'])
    s.shapes.add_picture(img_rev, Inches(3.0), Inches(1.1), Inches(10.0), Inches(2.9))

    # 曜日別棒グラフ（KPI の下）
    img_dow = chart_dow_bar(data['dow'])
    s.shapes.add_picture(img_dow, Inches(0.3), Inches(2.5), Inches(2.5), Inches(1.5))

    # アイテム別収益テーブル（下半分）
    txt(s, 'アイテム別収益 TOP7',
        Inches(0.3), Inches(4.2), Inches(4.0), Inches(0.35),
        size=10, bold=True, color=C_NAVY)

    rows = []
    for i, r in enumerate(data['items'][:7], 1):
        rows.append([
            str(i),
            r.get('itemName', '—')[:50],
            fmt_n(r.get('itemsPurchased', 0)),
            frev(r.get('itemRevenue', 0)),
        ])
    if rows:
        table(s,
              ['順位', 'アイテム名', '購入数', 'アイテム収益'],
              rows,
              Inches(0.3), Inches(4.6), Inches(12.7), Inches(2.75),
              font_size=9,
              col_widths=[Inches(0.5), Inches(8.5), Inches(1.7), Inches(2.0)])
    else:
        txt(s, 'eコマース収益データなし（設定または購入なし）',
            Inches(0.3), Inches(5.0), Inches(12.7), Inches(0.6),
            size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_spike_analysis(prs, site_name, data, year, month):
    """アクセス急増日 分析スライド"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'アクセス急増日 分析 — {site_name}  {year}年{month}月')

    spike_percentile = 80
    spikes, baseline = detect_spikes(data['device_daily'], percentile=spike_percentile)

    if not spikes:
        rect(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.0),
             C_COMMENT, C_COMMENT_BORDER)
        txt(s, '当月は急増日は検出されませんでした',
            Inches(1.7), Inches(2.9), Inches(9.9), Inches(0.7),
            size=14, color=RGBColor(0x2A, 0x5F, 0x8A), align=PP_ALIGN.CENTER)
        txt(s, f'（判定基準: 月中央値 {fmt_n(baseline)} セッション/日 の上位{100 - spike_percentile}%（{spike_percentile}パーセンタイル以上））',
            Inches(1.7), Inches(3.6), Inches(9.9), Inches(0.5),
            size=10, color=C_GRAY, align=PP_ALIGN.CENTER)
        return

    MEDIUM_JP = {
        'organic': 'オーガニック', 'referral': '参照', 'email': 'メール',
        'cpc': '有料CPC', 'social': 'SNS', 'none': 'ダイレクト',
        '(none)': 'ダイレクト', 'affiliate': 'アフィリエイト',
    }

    channel_daily = data.get('channel_daily', [])
    source_daily  = data.get('source_daily', [])

    # 参照元別・日次の平均（スパイク日を除く）を事前計算
    all_spike_dates = set(sp['date'] for sp in spikes)
    sm_by_date = {}
    all_sd_dates = set()
    for r in source_daily:
        d  = r.get('date', '')
        sm = r.get('sessionSourceMedium', '')
        sm_by_date.setdefault(sm, {})[d] = float(r.get('sessions', 0))
        all_sd_dates.add(d)
    non_spike_dates = all_sd_dates - all_spike_dates
    n_base = len(non_spike_dates) if non_spike_dates else 1
    sm_avg = {
        sm: sum(v for d, v in dv.items() if d in non_spike_dates) / n_base
        for sm, dv in sm_by_date.items()
    }

    show_spikes   = spikes[:3]
    box_h = Inches(1.95)
    gap   = Inches(0.1)
    top   = Inches(1.1)

    for i, sp in enumerate(show_spikes):
        y      = top + i * (box_h + gap)
        d_str  = sp['date']
        if len(d_str) == 8:
            d_label = f"{int(d_str[4:6])}月{int(d_str[6:8])}日"
        else:
            d_label = f"{int(d_str[5:7])}月{int(d_str[8:10])}日"

        analysis = analyze_spike_causes(sp['date'], channel_daily)
        comment  = build_spike_comment(analysis)

        # スパイク日の参照元別セッションを集計
        day_sources = {}
        for r in source_daily:
            if r.get('date', '') == sp['date']:
                sm   = r.get('sessionSourceMedium', '/ ')
                sess = float(r.get('sessions', 0))
                day_sources[sm] = day_sources.get(sm, 0) + sess
        top_sources = sorted(day_sources.items(), key=lambda x: x[1], reverse=True)[:5]

        # 左: 日付・倍率ボックス
        rect(s, Inches(0.3), y, Inches(2.9), box_h, C_NAVY)
        txt(s, d_label,
            Inches(0.3), y + Inches(0.1), Inches(2.9), Inches(0.52),
            size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, f'{fmt_n(sp["sessions"])} セッション',
            Inches(0.3), y + Inches(0.62), Inches(2.9), Inches(0.38),
            size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, f'通常の {sp["ratio"]:.1f} 倍',
            Inches(0.3), y + Inches(1.02), Inches(2.9), Inches(0.48),
            size=15, bold=True, color=RGBColor(0xF0, 0x78, 0x55), align=PP_ALIGN.CENTER)

        # 右: 原因分析ボックス
        rect(s, Inches(3.35), y, Inches(9.8), box_h, C_WHITE)

        # チャネル原因候補
        txt(s, '【原因候補】',
            Inches(3.5), y + Inches(0.06), Inches(2.5), Inches(0.28),
            size=9, bold=True, color=C_NAVY)
        txt(s, comment,
            Inches(3.5), y + Inches(0.34), Inches(9.5), Inches(0.38),
            size=11, bold=True, color=C_BLACK)

        # 参照元詳細
        txt(s, '【参照元 詳細 TOP5】',
            Inches(3.5), y + Inches(0.76), Inches(3.5), Inches(0.28),
            size=9, bold=True, color=C_NAVY)
        if top_sources:
            src_lines = []
            for sm, sess in top_sources:
                parts  = sm.split(' / ', 1)
                source = parts[0] if parts else sm
                medium = parts[1] if len(parts) > 1 else ''
                med_jp = MEDIUM_JP.get(medium.lower(), medium)
                avg    = sm_avg.get(sm, 0)
                delta  = sess - avg
                if avg > 0:
                    sign = '+' if delta >= 0 else ''
                    comp = f"（平均比 {sign}{fmt_n(delta)}件）"
                else:
                    comp = '（新規流入）'
                src_lines.append(f"・{source}（{med_jp}）: {fmt_n(sess)}件 {comp}")
            txt(s, '\n'.join(src_lines),
                Inches(3.5), y + Inches(1.04), Inches(9.5), Inches(0.82),
                size=9, color=C_GRAY)
        else:
            txt(s, 'データなし',
                Inches(3.5), y + Inches(1.04), Inches(9.5), Inches(0.3),
                size=9, color=C_GRAY)

    txt(s, f'判定基準: 月中央値 {fmt_n(baseline)} セッション/日 の上位{100 - spike_percentile}%（{spike_percentile}パーセンタイル以上）を急増日と判定',
        Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3),
        size=8, color=C_GRAY, italic=True)


def chart_funnel(funnel: dict):
    """購買ファネルの横棒グラフ"""
    steps  = ['view_item', 'add_to_cart', 'begin_checkout', 'purchase']
    labels = ['商品閲覧\nview_item', 'カート追加\nadd_to_cart',
              'チェックアウト\nbegin_checkout', '購入完了\npurchase']
    colors = ['#2272B5', '#5096C8', '#F07855', '#3BAD9B']
    counts = [funnel.get(e, 0) for e in steps]
    base   = counts[0] if counts[0] > 0 else 1
    pcts   = [c / base * 100 for c in counts]

    fig, ax = plt.subplots(figsize=(8.5, 4.2))
    y_pos = list(range(len(labels) - 1, -1, -1))
    for i, (y, pct, color) in enumerate(zip(y_pos, pcts, colors)):
        # 前段より計測数が多い等の異常値でも軸外にバー/ラベルが伸びて
        # savefig(bbox_inches='tight') が巨大画像を生成しないよう表示位置は上限で丸める
        bar_pct = min(pct, 130)
        ax.barh(y, bar_pct, color=color, alpha=0.85, height=0.6)
        count = counts[i]
        ax.text(min(pct + 1, 128), y, f'{count:,}件  ({pct:.1f}%)',
                va='center', fontsize=10, fontweight='bold', color='#333333', clip_on=True)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=9)
    ax.set_xlim(0, 130)
    ax.set_xlabel('商品閲覧を100%とした割合', fontsize=9)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{x:.0f}%'))
    ax.grid(axis='x', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('購買ファネル（GA4 イベント数）', fontsize=10, fontweight='bold', pad=8)
    fig.tight_layout()
    return chart_to_image(fig)


def slide_purchase_funnel(prs, site_name, data, year, month, period_label=None):
    """購買ファネル分析スライド"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'購買ファネル分析 — {site_name}  {period_label or f"{year}年{month}月"}')

    funnel = data.get('funnel', {})
    steps  = ['view_item', 'add_to_cart', 'begin_checkout', 'purchase']
    counts = [funnel.get(e, 0) for e in steps]
    base   = counts[0] if counts[0] > 0 else 1

    if base == 0:
        rect(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.0),
             C_COMMENT, C_COMMENT_BORDER)
        txt(s, 'eコマースイベントが計測されていません',
            Inches(1.7), Inches(2.9), Inches(9.9), Inches(0.65),
            size=15, bold=True, color=RGBColor(0x2A, 0x5F, 0x8A), align=PP_ALIGN.CENTER)
        txt(s, 'GA4 に view_item / add_to_cart / begin_checkout / purchase イベントが設定されているか確認してください。',
            Inches(1.7), Inches(3.6), Inches(9.9), Inches(0.5),
            size=10, color=RGBColor(0x2A, 0x5F, 0x8A), align=PP_ALIGN.CENTER)
        return

    # 左: ファネルチャート
    img = chart_funnel(funnel)
    s.shapes.add_picture(img, Inches(0.3), Inches(1.1), Inches(8.5), Inches(4.5))

    # 右: 転換率サマリー
    rx = Inches(9.1); rw = Inches(4.0)
    rect(s, rx, Inches(1.1), rw, Inches(5.6), C_WHITE)
    txt(s, '転換率サマリー', rx + Inches(0.2), Inches(1.2),
        rw - Inches(0.3), Inches(0.4), size=12, bold=True, color=C_NAVY)

    kpis = [
        ('商品閲覧 → カート追加',
         f'{counts[1] / max(counts[0], 1) * 100:.1f}%', C_BLUE),
        ('カート → チェックアウト',
         f'{counts[2] / max(counts[1], 1) * 100:.1f}%', C_ORANGE),
        ('チェックアウト → 購入',
         f'{counts[3] / max(counts[2], 1) * 100:.1f}%', C_GREEN),
        ('最終CVR（閲覧→購入）',
         f'{counts[3] / max(counts[0], 1) * 100:.2f}%', C_NAVY),
    ]
    for i, (label, val, col) in enumerate(kpis):
        ty = Inches(1.75) + i * Inches(1.2)
        rect(s, rx + Inches(0.15), ty, rw - Inches(0.3), Inches(1.0), C_LIGHT)
        rect(s, rx + Inches(0.15), ty, Inches(0.06), Inches(1.0), col)
        txt(s, label, rx + Inches(0.35), ty + Inches(0.1),
            rw - Inches(0.55), Inches(0.38), size=9, bold=True, color=C_BLACK)
        txt(s, val, rx + Inches(0.35), ty + Inches(0.48),
            rw - Inches(0.55), Inches(0.45), size=20, bold=True, color=col)

    final_cvr = counts[3] / max(counts[0], 1) * 100
    comment_box(s, f'最終CVR（商品閲覧→購入完了）{final_cvr:.2f}%。カート追加率{counts[1]/max(counts[0],1)*100:.1f}%。',
                Inches(6.55))


def chart_instagram_daily(insta_daily):
    """Instagram 日次セッション棒グラフ"""
    if not insta_daily:
        return _no_data_fig(8.5, 4.5)
    dates  = sorted(insta_daily.keys())
    values = [insta_daily[d] for d in dates]

    def d_label(d):
        return f"{int(d[4:6])}/{int(d[6:8])}" if len(d) == 8 else f"{int(d[5:7])}/{int(d[8:10])}"

    fig, ax = plt.subplots(figsize=(8.5, 4.5))
    bars = ax.bar(range(len(dates)), values, color='#F07855', alpha=0.85, edgecolor='white')

    # 上位3日にラベル
    top3_idx = sorted(range(len(values)), key=lambda i: values[i], reverse=True)[:3]
    for idx in top3_idx:
        ax.text(idx, values[idx] + max(values) * 0.02,
                f'{int(values[idx]):,}',
                ha='center', va='bottom', fontsize=8, color='#C0512A', fontweight='bold')

    step = max(1, len(dates) // 10)
    ax.set_xticks(range(0, len(dates), step))
    ax.set_xticklabels([d_label(dates[i]) for i in range(0, len(dates), step)], fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'))
    ax.set_ylabel('セッション数', fontsize=9)
    ax.grid(axis='y', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('Instagram 日次セッション数', fontsize=10, fontweight='bold', pad=6)
    fig.tight_layout()
    return chart_to_image(fig)


def slide_instagram(prs, site_name, data, year, month, period_label=None):
    """Instagram 流入分析：日次グラフ＋日別ランキング"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'Instagram 流入分析 — {site_name}  {period_label or f"{year}年{month}月"}')

    def is_instagram(sm):
        sm_l = sm.lower()
        return any(kw in sm_l for kw in ['instagram', 'ig / ', 'ig/', 'igshopping'])

    # Instagram 日次集計（全 Instagram 参照元を合算）
    insta_daily = {}
    for r in data.get('source_daily', []):
        if is_instagram(r.get('sessionSourceMedium', '')):
            d = r.get('date', '')
            insta_daily[d] = insta_daily.get(d, 0) + float(r.get('sessions', 0))

    # 日次グラフ（左）
    img = chart_instagram_daily(insta_daily)
    s.shapes.add_picture(img, Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.9))

    # 日別ランキング（右）
    txt(s, '日別ランキング',
        Inches(9.1), Inches(1.1), Inches(4.0), Inches(0.4),
        size=11, bold=True, color=C_NAVY)

    ranked = sorted(insta_daily.items(), key=lambda x: x[1], reverse=True)[:15]
    rows = []
    for rank, (d, sess) in enumerate(ranked, 1):
        if len(d) == 8:
            d_lbl = f"{int(d[4:6])}月{int(d[6:8])}日"
        else:
            d_lbl = f"{int(d[5:7])}月{int(d[8:10])}日"
        rows.append([str(rank), d_lbl, fmt_n(sess)])

    if rows:
        table(s, ['順位', '日付', 'セッション'],
              rows,
              Inches(9.1), Inches(1.55), Inches(4.0), Inches(5.3),
              font_size=9)
    else:
        txt(s, 'データなし',
            Inches(9.1), Inches(3.5), Inches(4.0), Inches(0.5),
            size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    txt(s, '※ ig/social・l.instagram.com・IGShopping・instagram.com の合算',
        Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
        size=8, color=C_GRAY, italic=True)


def slide_monthly_trend(prs, site_name, data):
    """直近3ヶ月推移"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'直近3ヶ月推移 — {site_name}')
    img = chart_monthly_trend(data['monthly_trend'])
    s.shapes.add_picture(img, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.9))


def slide_benchmark_summary(prs, site_name, benchmark, year, month, period_label=None):
    """他サイトと比較して目立って低い指標のコメント（最終ページ）"""
    s = blank(prs)
    bg(s, C_LIGHT)
    header(s, f'サイト間比較コメント — {site_name}  {period_label or f"{year}年{month}月"}')

    comments = benchmark_lowlights(site_name, benchmark)
    rect(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.2), C_WHITE)

    if comments:
        txt(s, '他サイトと比較して目立って低い指標', Inches(1.1), Inches(1.55), Inches(11.1), Inches(0.5),
            size=15, bold=True, color=C_RED)
        txt(s, '\n'.join(comments), Inches(1.1), Inches(2.2), Inches(11.1), Inches(3.5),
            size=13, color=C_GRAY)
    else:
        txt(s, '他サイトと比較して顕著に低い指標はありません', Inches(1.1), Inches(3.5), Inches(11.1), Inches(0.6),
            size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    txt(s, '※ CVR・エンゲージメント率について、データのある他サイトの中央値と比較（中央値の70%未満を「目立って低い」と判定）',
        Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), size=8, color=C_GRAY, italic=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_detail_report(site_name, property_id, year, month, output_dir,
                        search_console_url=None, has_ecommerce=False, benchmark=None):
    start, end = get_date_range(year, month)

    print(f'[{site_name}] GA4 データ取得中...')
    data = fetch_detail_data(property_id, year, month)

    gsc_rows = None
    if search_console_url:
        print(f'[{site_name}] Search Console データ取得中...')
        gsc_rows = fetch_search_console(search_console_url, start, end, limit=10)
        print(f'  取得件数: {len(gsc_rows)} 件')

    prs = new_prs()
    print(f'[{site_name}] スライド生成中...')

    slide_cover(prs, site_name, year, month)
    slide_traffic_summary(prs, site_name, data, year, month)
    if has_ecommerce:
        slide_purchase_funnel(prs, site_name, data, year, month)
    slide_demographics(prs, site_name, data, year, month)
    slide_page_traffic(prs, site_name, data, year, month)
    slide_source_medium(prs, site_name, data, year, month)
    slide_campaign(prs, site_name, data, year, month)
    slide_search_query(prs, site_name, year, month, gsc_rows=gsc_rows)
    slide_revenue_combined(prs, site_name, data, year, month)
    slide_monthly_trend(prs, site_name, data)
    slide_instagram(prs, site_name, data, year, month)
    slide_spike_analysis(prs, site_name, data, year, month)
    if benchmark:
        slide_benchmark_summary(prs, site_name, benchmark, year, month)

    fname = f'{site_name}_詳細レポート_{year}{month:02d}.pptx'
    out_path = os.path.join(output_dir, fname)
    prs.save(out_path)
    print(f'保存完了: {out_path}')
    return out_path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 複数月まとめレポート
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def fetch_quarterly_data(property_id, start, end, months_list):
    """指定期間のデータ一括取得（quarterly report用）"""
    pid = property_id

    # 比較期間: 同じ月数分だけ前
    n = len(months_list)
    py, pm = months_list[0]
    for _ in range(n):
        py, pm = prev_month(py, pm)
    prev_start = get_date_range(py, pm)[0]
    prev_end_y, prev_end_m = months_list[0]
    prev_end_y, prev_end_m = prev_month(prev_end_y, prev_end_m)
    prev_end = get_date_range(prev_end_y, prev_end_m)[1]

    print('  基本集計...')
    t_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    totals = t_rows[0] if t_rows else {}

    p_rows = safe_run(pid, prev_start, prev_end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    prev_totals = p_rows[0] if p_rows else {}

    # 前年同期: 対象月リストをそのまま1年前にシフト
    yoy_start = get_date_range(months_list[0][0] - 1, months_list[0][1])[0]
    yoy_end   = get_date_range(months_list[-1][0] - 1, months_list[-1][1])[1]
    y_rows = safe_run(pid, yoy_start, yoy_end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'])
    yoy_totals = y_rows[0] if y_rows else {}

    print('  デバイス別日次...')
    device_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'deviceCategory'], limit=300)
    device_daily.sort(key=lambda r: r.get('date', ''))

    print('  ユーザー属性...')
    gender_rows    = safe_run(pid, start, end, ['sessions'], ['userGender'], limit=5)
    age_rows       = safe_run(pid, start, end, ['sessions'], ['userAgeBracket'], limit=10)
    retention_rows = safe_run(pid, start, end, ['sessions'], ['newVsReturning'], limit=5)

    print('  ページタイトル...')
    page_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions'],
        ['pageTitle'], order='sessions', limit=10)

    print('  参照元/メディア...')
    source_rows = safe_run(pid, start, end,
        ['sessions', 'totalUsers', 'newUsers', 'engagedSessions', 'totalRevenue'],
        ['sessionSourceMedium'], order='sessions', limit=10)

    print('  キャンペーン...')
    campaign_rows = safe_run(pid, start, end,
        ['sessions', 'totalRevenue'],
        ['sessionCampaignName', 'sessionSourceMedium'], order='sessions', limit=10)

    print('  日次収益...')
    revenue_daily = safe_run(pid, start, end,
        ['totalRevenue', 'sessions'], ['date'], limit=100)
    revenue_daily.sort(key=lambda r: r.get('date', ''))

    print('  アイテム別収益...')
    item_rows = safe_run(pid, start, end,
        ['itemsPurchased', 'itemRevenue'], ['itemName'], order='itemRevenue', limit=10)

    print('  チャネル別日次...')
    channel_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'sessionDefaultChannelGroup'], limit=1500)

    print('  参照元別日次...')
    source_daily = safe_run(pid, start, end,
        ['sessions'], ['date', 'sessionSourceMedium'], limit=3000)

    print('  曜日別...')
    dow_rows = safe_run(pid, start, end,
        ['sessions', 'totalRevenue'], ['dayOfWeek'], limit=7)
    dow_rows.sort(key=lambda r: int(r.get('dayOfWeek', 0)))

    print('  月別推移...')
    monthly_trend = []
    for my, mm in months_list:
        ms, me = get_date_range(my, mm)
        rows = safe_run(pid, ms, me, ['sessions', 'totalRevenue'], limit=1)
        r = rows[0] if rows else {}
        monthly_trend.append({
            'year': my, 'month': mm,
            'sessions': float(r.get('sessions', 0)),
            'revenue':  float(r.get('totalRevenue', 0)),
        })

    return {
        'totals':        totals,
        'prev_totals':   prev_totals,
        'yoy_totals':    yoy_totals,
        'device_daily':  device_daily,
        'gender':        gender_rows,
        'age':           age_rows,
        'retention':     retention_rows,
        'pages':         page_rows,
        'sources':       source_rows,
        'campaigns':     campaign_rows,
        'revenue_daily': revenue_daily,
        'items':         item_rows,
        'dow':           dow_rows,
        'monthly_trend': monthly_trend,
        'channel_daily': channel_daily,
        'source_daily':  source_daily,
    }


def slide_cover_quarterly(prs, site_name, start_ym, end_ym):
    """複数月まとめレポート 表紙"""
    s = blank(prs)
    bg(s, C_WHITE)
    rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), C_NAVY)
    rect(s, Inches(0.5), Inches(1.8), Inches(0.08), Inches(3.5), C_NAVY)
    txt(s, site_name,
        Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.1),
        size=28, bold=True, color=C_NAVY)
    txt(s, 'GA4 詳細アクセス・収益レポート（複数月まとめ）',
        Inches(0.8), Inches(2.9), Inches(11), Inches(0.8),
        size=20, color=C_BLUE)
    sy, sm = start_ym
    ey, em = end_ym
    txt(s, f'{sy}年{sm}月 〜 {ey}年{em}月',
        Inches(0.8), Inches(3.7), Inches(11), Inches(0.7),
        size=18, color=C_BLACK)
    txt(s, '※ スパイク分析のみ月別。その他スライドは期間全体の集計値です。',
        Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
        size=11, color=C_GRAY)
    txt(s, 'Powered by Google Analytics 4',
        Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
        size=11, color=C_GRAY)


def build_quarterly_report(site_name, property_id, year, month, output_dir,
                           n_months=3, search_console_url=None, has_ecommerce=False):
    """直近 n_months ヶ月まとめレポート（スパイク分析のみ月別スライド）"""
    # 対象月リスト（古い順）
    months_list = []
    cy, cm = year, month
    for _ in range(n_months):
        months_list.insert(0, (cy, cm))
        cy, cm = prev_month(cy, cm)

    start = get_date_range(months_list[0][0], months_list[0][1])[0]
    end   = get_date_range(year, month)[1]
    sy, sm = months_list[0]
    period_label = f'{sy}年{sm}月〜{year}年{month}月'

    print(f'[{site_name}] GA4 データ取得中（{period_label}）...')
    data = fetch_quarterly_data(property_id, start, end, months_list)

    if has_ecommerce:
        print(f'  購買ファネル...')
        try:
            data['funnel'] = fetch_funnel_data(property_id, start, end)
        except Exception as ex:
            print(f'  [skip] 購買ファネル: {ex}')
            data['funnel'] = {'view_item': 0, 'add_to_cart': 0,
                              'begin_checkout': 0, 'purchase': 0}

    gsc_rows = None
    if search_console_url:
        print(f'[{site_name}] Search Console データ取得中...')
        gsc_rows = fetch_search_console(search_console_url, start, end, limit=10)
        print(f'  取得件数: {len(gsc_rows)} 件')

    prs = new_prs()
    print(f'[{site_name}] スライド生成中...')

    slide_cover_quarterly(prs, site_name, months_list[0], (year, month))
    slide_traffic_summary(prs, site_name, data, year, month, period_label=period_label)
    if has_ecommerce:
        slide_purchase_funnel(prs, site_name, data, year, month,
                              period_label=period_label)
    slide_demographics(prs, site_name, data, year, month, period_label=period_label)
    slide_page_traffic(prs, site_name, data, year, month, period_label=period_label)
    slide_source_medium(prs, site_name, data, year, month, period_label=period_label)
    slide_campaign(prs, site_name, data, year, month, period_label=period_label)
    slide_search_query(prs, site_name, year, month, gsc_rows=gsc_rows,
                       period_label=period_label)
    slide_revenue_combined(prs, site_name, data, year, month, period_label=period_label)
    slide_monthly_trend(prs, site_name, data)
    slide_instagram(prs, site_name, data, year, month, period_label=period_label)

    # スパイク分析：月別（quarterly data から対象月の日次データをフィルタ）
    # GA4の日付はYYYYMMDD形式のためハイフンなし形式で比較
    for my, mm in months_list:
        ms, me = get_date_range(my, mm)
        ms_raw = ms.replace('-', '')   # '2026-03-01' → '20260301'
        me_raw = me.replace('-', '')   # '2026-03-31' → '20260331'
        month_data = {
            'device_daily':  [r for r in data['device_daily']
                               if ms_raw <= r.get('date', '') <= me_raw],
            'channel_daily': [r for r in data['channel_daily']
                               if ms_raw <= r.get('date', '') <= me_raw],
            'source_daily':  [r for r in data['source_daily']
                               if ms_raw <= r.get('date', '') <= me_raw],
        }
        slide_spike_analysis(prs, site_name, month_data, my, mm)

    fname = f'{site_name}_詳細レポート_{sy}{sm:02d}-{year}{month:02d}.pptx'
    out_path = os.path.join(output_dir, fname)
    prs.save(out_path)
    print(f'保存完了: {out_path}')
    return out_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--site',      default='cotopaxiオフィシャルサイト')
    parser.add_argument('--month',     default=None, help='YYYY-MM（省略で前月）')
    parser.add_argument('--quarterly', action='store_true',
                        help='複数月まとめレポートを生成（--months で月数指定）')
    parser.add_argument('--months',    type=int, default=3,
                        help='まとめる月数（デフォルト3）')
    args = parser.parse_args()

    with open(CONFIG_YML, encoding='utf-8') as f:
        config = yaml.safe_load(f)

    if args.month:
        year, month = int(args.month[:4]), int(args.month[5:7])
    else:
        today = date.today()
        year, month = prev_month(today.year, today.month)

    needle = args.site.lower()
    matched = [s for s in config['sites']
               if s['name'].lower() == needle or needle in s['name'].lower()]
    if not matched:
        print(f'サイトが見つかりません: {args.site}')
        for s in config['sites']:
            print(f'  {s["name"]}')
        sys.exit(1)

    site = matched[0]
    output_dir = config.get('output_dir', '.')
    has_ecommerce = site.get('has_ecommerce', False)
    if args.quarterly:
        build_quarterly_report(
            site['name'], site['property_id'], year, month, output_dir,
            n_months=args.months,
            search_console_url=site.get('search_console_url'),
            has_ecommerce=has_ecommerce,
        )
    else:
        print('ベンチマークデータ取得中（他サイト比較用）...')
        benchmark = fetch_benchmark_metrics(config['sites'], year, month)
        build_detail_report(
            site['name'], site['property_id'], year, month, output_dir,
            search_console_url=site.get('search_console_url'),
            has_ecommerce=has_ecommerce,
            benchmark=benchmark,
        )


if __name__ == '__main__':
    main()
