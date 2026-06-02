"""
EC 昨対比較分析 コアモジュール
load() → compute_stats() → make_pos() / make_int() の順で呼び出す
"""
import io
from types import SimpleNamespace
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 日本語フォント（リポジトリ同梱の ipaexg.ttf を直接指定）──
import pathlib as _pl
_FONT_PATH = str(_pl.Path(__file__).parent / 'fonts' / 'ipaexg.ttf')
if _pl.Path(_FONT_PATH).exists():
    fm.fontManager.addfont(_FONT_PATH)
    plt.rcParams['font.family'] = fm.FontProperties(fname=_FONT_PATH).get_name()
else:
    # ローカル開発環境フォールバック（Windows）
    for _c in ['Yu Gothic', 'Meiryo', 'MS Gothic']:
        _hits = [f for f in fm.fontManager.ttflist if _c.lower() in f.name.lower()]
        if _hits:
            plt.rcParams['font.family'] = _hits[0].name
            break
plt.rcParams['axes.unicode_minus'] = False

C_TEAL='#0DB4C6'; C_DARK='#1A2535'; C_GOLD='#F5A623'; C_GREEN='#27AE60'
C_RED='#E74C3C'; C_PURPLE='#8E44AD'; C_LGRAY='#F5F5F5'; C_GRAY='#9E9E9E'
C_WHITE='#FFFFFF'; C_NAVY='#1A376C'

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ══════════════════════════════════════
# データ読み込み
# ══════════════════════════════════════
def load(file_obj):
    """
    ファイルパス(str)またはファイルオブジェクトを受け取り DataFrame を返す。
    ヘッダーは5行目（header=4）、列順は固定。
    """
    df = pd.read_excel(file_obj, header=4, engine='calamine')
    df.columns = ['実績年月','JAN','ブランドCD','ブランド名','商品CD','商品名',
                  'カラーCD','カラー','サイズCD','サイズ','請求先','得意先',
                  '担当者','売上数','売上金額','旧在庫区分','最新在庫区分','詳細分類']
    df['実績年月'] = pd.to_datetime(df['実績年月'], errors='coerce')
    df = df.dropna(subset=['実績年月'])
    for c in ['売上数','売上金額']:
        df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)
    df['isD'] = df['旧在庫区分'] == 'D'
    df['月'] = df['実績年月'].dt.month.astype(str) + '月'
    # 店舗名の統一（名称変更で別店舗に見えるケースを統合）
    df['得意先'] = df['得意先'].replace({'Fav_Our_Planet': 'ALCO online'})
    return df

# ══════════════════════════════════════
# 統計計算
# ══════════════════════════════════════
def compute_stats(d25, d26):
    s = SimpleNamespace()
    s.d25 = d25; s.d26 = d26

    # 月・年・ブランドを自動検出
    months_ord = sorted(d25['実績年月'].dt.month.unique())
    s.MONTHS = [f'{m}月' for m in months_ord]
    s.brand  = d25['ブランド名'].dropna().mode().iloc[0] if d25['ブランド名'].dropna().any() else 'EC'
    s.year25 = int(d25['実績年月'].dt.year.mode().iloc[0])
    s.year26 = int(d26['実績年月'].dt.year.mode().iloc[0])
    s.period = f'{s.MONTHS[0]}〜{s.MONTHS[-1]}' if len(s.MONTHS) > 1 else s.MONTHS[0]

    # プロパー品
    s.p25 = d25[~d25['isD']].copy()
    s.p26 = d26[~d26['isD']].copy()

    def monthly(df):
        return df.groupby('月').agg(数量=('売上数','sum'), 金額=('売上金額','sum')
                ).reindex(s.MONTHS, fill_value=0)

    def safe_yoy(a, b):
        return (b/a-1)*100 if a > 0 else 0.0

    ma25 = monthly(d25); ma26 = monthly(d26)
    s.UNIT_G_A = [safe_yoy(ma25.loc[m,'数量'], ma26.loc[m,'数量']) for m in s.MONTHS]
    s.REV_G_A  = [safe_yoy(ma25.loc[m,'金額'],  ma26.loc[m,'金額'])  for m in s.MONTHS]
    s.TOT_UNIT_A = safe_yoy(d25['売上数'].sum(),   d26['売上数'].sum())
    s.TOT_REV_A  = safe_yoy(d25['売上金額'].sum(), d26['売上金額'].sum())
    s.UNIT_25_MO = [ma25.loc[m,'数量'] for m in s.MONTHS]
    s.UNIT_26_MO = [ma26.loc[m,'数量'] for m in s.MONTHS]
    s.REV_25_MO  = [ma25.loc[m,'金額']  for m in s.MONTHS]
    s.REV_26_MO  = [ma26.loc[m,'金額']  for m in s.MONTHS]

    ms25 = monthly(s.p25); ms26 = monthly(s.p26)
    s.UNIT_G_P = [safe_yoy(ms25.loc[m,'数量'], ms26.loc[m,'数量']) for m in s.MONTHS]
    s.REV_G_P  = [safe_yoy(ms25.loc[m,'金額'],  ms26.loc[m,'金額'])  for m in s.MONTHS]
    s.TOT_UNIT_P = safe_yoy(s.p25['売上数'].sum(),   s.p26['売上数'].sum())
    s.TOT_REV_P  = safe_yoy(s.p25['売上金額'].sum(), s.p26['売上金額'].sum())

    def top_prods(df26, df25, n=10):
        t = df26.groupby('商品名')['売上数'].sum().sort_values(ascending=False).head(n)
        total = df26['売上数'].sum()
        rank25 = {nm: i+1 for i,nm in enumerate(
            df25.groupby('商品名')['売上数'].sum().sort_values(ascending=False).index)}
        return [(nm, v/total*100 if total>0 else 0,
                 nm not in rank25 or rank25[nm]>10) for nm,v in t.items()]

    def top_colors_fn(df26, df25, n=10):
        t = df26.groupby('カラー')['売上数'].sum().sort_values(ascending=False)
        total = df26['売上数'].sum()
        c25 = set(df25['カラー'].unique())
        result = [(nm, v/total*100 if total>0 else 0, nm not in c25) for nm,v in t.head(n).items()]
        result.append(('その他', (total-t.head(n).sum())/total*100 if total>0 else 0, False))
        return result, len(set(df26['カラー'].unique()) - c25)

    def size_dist(df):
        t = df.groupby('サイズ')['売上数'].sum().sort_values(ascending=False)
        total = df['売上数'].sum()
        return {nm: v/total*100 if total>0 else 0 for nm,v in t.items()}

    _nm_map = {'ｻﾝﾀﾞﾙ':'サンダル','ｸﾛｯｸﾞ':'クロッグ','ﾄﾝｸﾞ':'トング',
               'ｼｬﾜｰ':'シャワー','ｽﾘｯﾎﾟﾝ':'スリッポン','ﾌﾞｰﾂ':'ブーツ'}
    def fix_nm(nm):
        for k,v in _nm_map.items(): nm = nm.replace(k,v)
        return nm
    def cat_dist(df):
        t = df.groupby('詳細分類')['売上数'].sum().sort_values(ascending=False)
        total = df['売上数'].sum()
        return [(fix_nm(nm), v/total*100 if total>0 else 0) for nm,v in t.items()]

    s.TOP10_A = top_prods(d26, d25)
    s.TOP_COLORS_A, s.NEW_COLORS_A = top_colors_fn(d26, d25)
    sz25a = size_dist(d25); sz26a = size_dist(d26)
    s.SZ_LABELS_A = [k for k in sz26a if sz26a[k]>=1.0][:8]
    s.SZ_2025_A = [sz25a.get(k,0) for k in s.SZ_LABELS_A]
    s.SZ_2026_A = [sz26a.get(k,0) for k in s.SZ_LABELS_A]
    s.CAT25_A = cat_dist(d25); s.CAT26_A = cat_dist(d26)

    s.TOP10_P = top_prods(s.p26, s.p25)
    s.TOP_COLORS_P, s.NEW_COLORS_P = top_colors_fn(s.p26, s.p25)
    sz25p = size_dist(s.p25); sz26p = size_dist(s.p26)
    s.SZ_LABELS_P = [k for k in sz26p if sz26p[k]>=1.0][:8]
    s.SZ_2025_P = [sz25p.get(k,0) for k in s.SZ_LABELS_P]
    s.SZ_2026_P = [sz26p.get(k,0) for k in s.SZ_LABELS_P]
    s.CAT25_P = cat_dist(s.p25); s.CAT26_P = cat_dist(s.p26)
    s.PORT = {
        str(s.year25): (s.p25['商品名'].nunique(), s.p25['JAN'].nunique(), s.p25['カラー'].nunique()),
        str(s.year26): (s.p26['商品名'].nunique(), s.p26['JAN'].nunique(), s.p26['カラー'].nunique()),
    }

    # 在庫区分Dアイテム
    s.d26_sale    = d26[d26['isD']].copy()
    s.sale_qty_25 = d25[d25['isD']]['売上数'].sum()
    s.sale_qty_26 = s.d26_sale['売上数'].sum()
    s.sale_ratio_26 = s.sale_qty_26 / d26['売上数'].sum() * 100 if d26['売上数'].sum() > 0 else 0
    s.sale_prods26  = s.d26_sale.groupby('商品名').agg(
        数量=('売上数','sum'), 金額=('売上金額','sum')).sort_values('数量',ascending=False).head(8)

    # サイト別
    sg25 = d25.groupby('得意先').agg(数量=('売上数','sum'), 金額=('売上金額','sum'))
    sg26 = d26.groupby('得意先').agg(数量=('売上数','sum'), 金額=('売上金額','sum'))
    s.site_new  = '・'.join(si for si in sg26.index if si not in sg25.index)
    s.site_gone = '・'.join(si for si in sg25.index if si not in sg26.index)
    common = [si for si in sg26.index if si in sg25.index and sg25.loc[si,'数量']>0]
    yoy_qty = {si: (sg26.loc[si,'数量']/sg25.loc[si,'数量']-1)*100 for si in common}
    s.worst_site = min(yoy_qty, key=yoy_qty.get) if yoy_qty else ''
    s.worst_yoy  = yoy_qty.get(s.worst_site, 0)

    return s

# ══════════════════════════════════════
# チャート生成
# ══════════════════════════════════════
def _to_buf(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    buf.seek(0); plt.close(fig); return buf

def ch_monthly(ug, rg, tot_u, tot_r, months):
    fig, ax = plt.subplots(figsize=(9.5, 5.2), facecolor='white')
    x = np.arange(len(months)); w = 0.35
    bu = ax.bar(x-w/2, ug, w, label='販売点数 前年比',
                color=[C_GREEN if v>=0 else C_RED for v in ug], alpha=0.85, zorder=3)
    br = ax.bar(x+w/2, rg, w, label='売上高 前年比',
                color=[C_TEAL if v>=0 else '#CCAAAA' for v in rg], alpha=0.85, zorder=3)
    for bar, v in list(zip(bu,ug))+list(zip(br,rg)):
        ofs = 1.8 if v>=0 else -3.5
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+ofs, f'{v:+.1f}%',
                ha='center', va='bottom' if v>=0 else 'top',
                fontsize=12, fontweight='bold', color=C_DARK)
    ax.axhline(0, color=C_DARK, lw=0.8, zorder=4)
    ax.set_xticks(x); ax.set_xticklabels(months, fontsize=14)
    ax.set_ylabel('前年比（%）', fontsize=12)
    ax.legend(fontsize=12, loc='upper left')
    ax.set_ylim(min(min(ug),min(rg))-14, max(max(ug),max(rg))+18)
    ax.grid(axis='y', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    u_s = f'+{tot_u:.1f}%' if tot_u>=0 else f'{tot_u:.1f}%'
    r_s = f'+{tot_r:.1f}%' if tot_r>=0 else f'{tot_r:.1f}%'
    ax.annotate(f'累計：販売点数 {u_s}  ／  売上高 {r_s}',
                xy=(0.5,0.02), xycoords='axes fraction', ha='center', fontsize=12,
                color=C_TEAL, fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.35', fc='#E8F9FA', ec=C_TEAL))
    fig.tight_layout(); return _to_buf(fig)

def ch_ranking(top10):
    names=[t[0] for t in top10]; shares=[t[1] for t in top10]; isnew=[t[2] for t in top10]
    fig, ax = plt.subplots(figsize=(9.5, 5.8), facecolor='white')
    y = np.arange(len(names))
    bars = ax.barh(y, shares, color=[C_GOLD if n else C_TEAL for n in isnew], alpha=0.85, height=0.65, zorder=3)
    ax.set_yticks(y)
    ax.set_yticklabels([f'  NEW  {nm}' if nw else f'       {nm}' for nm,nw in zip(names,isnew)], fontsize=12)
    for bar, v in zip(bars, shares):
        ax.text(bar.get_width()+0.2, bar.get_y()+bar.get_height()/2,
                f'{v:.1f}%', va='center', fontsize=12, fontweight='bold', color=C_DARK)
    ax.set_xlabel('販売構成比（%）', fontsize=12)
    ax.set_xlim(0, max(shares)+5); ax.invert_yaxis()
    ax.grid(axis='x', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    ax.legend(handles=[mpatches.Patch(color=C_TEAL, alpha=0.85, label='前年からの継続商品'),
                       mpatches.Patch(color=C_GOLD, alpha=0.85, label='新ライン・ランクアップ')],
              fontsize=11, loc='lower right')
    fig.tight_layout(); return _to_buf(fig)

def ch_color(top_colors, new_colors):
    names=[t[0] for t in top_colors]; shares=[t[1] for t in top_colors]; isnew=[t[2] for t in top_colors]
    fig, ax = plt.subplots(figsize=(9.5, 5.8), facecolor='white')
    y = np.arange(len(names))
    cols = [C_GOLD if n else (C_GRAY if nm=='その他' else C_TEAL) for nm,n in zip(names,isnew)]
    bars = ax.barh(y, shares, color=cols, alpha=0.85, height=0.65, zorder=3)
    ax.set_yticks(y)
    ax.set_yticklabels([f'  NEW  {nm}' if nw else f'       {nm}' for nm,nw in zip(names,isnew)], fontsize=12)
    for bar, v in zip(bars, shares):
        ax.text(bar.get_width()+0.3, bar.get_y()+bar.get_height()/2,
                f'{v:.1f}%', va='center', fontsize=12, color=C_DARK)
    ax.set_xlabel('構成比（%）', fontsize=12)
    ax.set_xlim(0, max(shares)+10); ax.invert_yaxis()
    ax.grid(axis='x', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    ax.legend(handles=[mpatches.Patch(color=C_TEAL, alpha=0.85, label='継続カラー'),
                       mpatches.Patch(color=C_GOLD, alpha=0.85, label=f'新カラー（計{new_colors}色追加）')],
              fontsize=11, loc='lower right')
    ax.annotate(f'新色 {new_colors}色 追加', xy=(0.97,0.02), xycoords='axes fraction',
                ha='right', fontsize=11, color=C_GOLD, fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.3', fc='#FFF8E8', ec=C_GOLD))
    fig.tight_layout(); return _to_buf(fig)

def ch_size(sz_labels, sz_2025, sz_2026):
    fig, ax = plt.subplots(figsize=(9.5, 5.0), facecolor='white')
    x = np.arange(len(sz_labels)); w = 0.35
    ax.bar(x-w/2, sz_2025, w, label='前年', color=C_GRAY, alpha=0.65, zorder=3)
    ax.bar(x+w/2, sz_2026, w, label='今年', color=C_TEAL, alpha=0.85, zorder=3)
    ax.set_xticks(x); ax.set_xticklabels(sz_labels, fontsize=11, rotation=30)
    ax.set_ylabel('構成比（%）', fontsize=12); ax.legend(fontsize=12)
    ax.set_ylim(0, max(max(sz_2025),max(sz_2026))+6)
    ax.grid(axis='y', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(); return _to_buf(fig)

def ch_category(cat25, cat26):
    all_cats = list(dict.fromkeys([t[0] for t in cat26]+[t[0] for t in cat25]))
    v25 = [next((v for n,v in cat25 if n==c),0) for c in all_cats]
    v26 = [next((v for n,v in cat26 if n==c),0) for c in all_cats]

    # 前年・今年どちらも5%未満のカテゴリを「その他」に束ねる
    main_cats, main_v25, main_v26 = [], [], []
    oth25 = oth26 = 0.0
    for c, a, b in zip(all_cats, v25, v26):
        if max(a, b) >= 5.0:
            main_cats.append(c); main_v25.append(a); main_v26.append(b)
        else:
            oth25 += a; oth26 += b
    if oth25 > 0 or oth26 > 0:
        main_cats.append('その他'); main_v25.append(oth25); main_v26.append(oth26)
    all_cats, v25, v26 = main_cats, main_v25, main_v26

    n = len(all_cats)
    rot  = 40 if n > 5 else 0
    fs   = max(9, 12 - max(0, n - 5))
    figh = 5.5 if rot else 4.5
    fig, ax = plt.subplots(figsize=(8.5, figh), facecolor='white')
    x = np.arange(n); w = 0.35
    ax.bar(x-w/2, v25, w, label='前年', color=C_GRAY, alpha=0.65, zorder=3)
    b2 = ax.bar(x+w/2, v26, w, label='今年', color=C_TEAL, alpha=0.85, zorder=3)
    # 前年ラベル（小さめで補足表示）
    for bar, v in zip(ax.patches[:n], v25):
        if v > 0:
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                    f'{v:.1f}%', ha='center', fontsize=fs-1, color=C_GRAY)
    # 今年ラベル（前年比diffは今年が8%以上の場合のみ表示して重なりを防ぐ）
    top = max(max(v25), max(v26))
    for bar, v, v0 in zip(b2, v26, v25):
        diff = v - v0
        show_diff = abs(diff) > 0.5 and v >= 8.0
        label = f'{v:.1f}%' + (f'\n({diff:+.1f}pt)' if show_diff else '')
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3,
                label, ha='center', fontsize=fs,
                color=C_GREEN if diff > 0.5 else (C_RED if diff < -0.5 else C_DARK),
                fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(all_cats, fontsize=fs+1, rotation=rot,
                       ha='right' if rot else 'center')
    ax.set_ylabel('構成比（%）', fontsize=12); ax.legend(fontsize=12)
    ax.set_ylim(0, top + 18)
    ax.grid(axis='y', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(); return _to_buf(fig)

def ch_portfolio(port, year25, year26):
    cats = ['取扱モデル数', 'SKU数', 'カラー数']
    v25 = list(port[str(year25)]); v26 = list(port[str(year26)])
    fig, ax = plt.subplots(figsize=(8.5, 4.5), facecolor='white')
    x = np.arange(len(cats)); w = 0.35
    b1 = ax.bar(x-w/2, v25, w, label=f'{year25}年', color=C_GRAY, alpha=0.65, zorder=3)
    b2 = ax.bar(x+w/2, v26, w, label=f'{year26}年', color=C_TEAL, alpha=0.85, zorder=3)
    for bar, v in zip(b1, v25):
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3, str(v),
                ha='center', fontsize=13, color=C_DARK)
    for bar, v, v0 in zip(b2, v26, v25):
        g = (v/v0-1)*100 if v0>0 else 0
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.3, f'{v}\n({g:+.0f}%)',
                ha='center', fontsize=12, color=C_TEAL if g>0 else C_RED, fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(cats, fontsize=13)
    ax.legend(fontsize=12); ax.set_ylim(0, max(v26)+20)
    ax.grid(axis='y', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(); return _to_buf(fig)

def _label_stacked(ax, x, prop_val, sale_val, fs_prop=13, fs_sale=12, thresh=6):
    ax.text(x, prop_val/2, f'{prop_val:.1f}%', ha='center', va='center',
            fontsize=fs_prop, fontweight='bold', color=C_WHITE)
    if sale_val > 0.3:
        if sale_val >= thresh:
            ax.text(x, prop_val+sale_val/2, f'{sale_val:.1f}%',
                    ha='center', va='center', fontsize=fs_sale, fontweight='bold', color=C_WHITE)
        else:
            ax.annotate(f'区分D: {sale_val:.1f}%',
                        xy=(x, prop_val+sale_val), xytext=(x, prop_val+sale_val+8),
                        ha='center', va='bottom', fontsize=fs_sale, fontweight='bold', color=C_RED,
                        arrowprops=dict(arrowstyle='->', color=C_RED, lw=1.2))

def ch_sale_ratio(s):
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 6.2), facecolor='white')
    sale_r_25 = s.d25[s.d25['isD']]['売上数'].sum() / s.d25['売上数'].sum() * 100 \
                if s.d25['売上数'].sum() > 0 else 0
    prop_r = [100-sale_r_25, 100-s.sale_ratio_26]
    sale_r = [sale_r_25, s.sale_ratio_26]
    x = np.arange(2); w = 0.5
    ax1.bar(x, prop_r, w, label='プロパー品', color=C_TEAL, alpha=0.85, zorder=3)
    ax1.bar(x, sale_r, w, bottom=prop_r, label='在庫区分Dアイテム', color=C_RED, alpha=0.75, zorder=3)
    for i,(p,sa) in enumerate(zip(prop_r, sale_r)):
        _label_stacked(ax1, i, p, sa)
    ax1.set_xticks(x); ax1.set_xticklabels([f'{s.year25}年', f'{s.year26}年'], fontsize=13)
    ax1.set_ylabel('構成比（%）', fontsize=12); ax1.set_ylim(0, 125); ax1.legend(fontsize=11)
    ax1.set_title('在庫区分別 構成比', fontsize=12)
    ax1.grid(axis='y', alpha=0.3, zorder=1); ax1.spines[['top','right']].set_visible(False)
    mo_total = (s.p26.groupby('月')['売上数'].sum()
                + s.d26_sale.groupby('月')['売上数'].sum().reindex(s.MONTHS, fill_value=0))
    mo_sale = s.d26_sale.groupby('月')['売上数'].sum().reindex(s.MONTHS, fill_value=0)
    mo_prop = s.p26.groupby('月')['売上数'].sum().reindex(s.MONTHS, fill_value=0)
    denom = mo_total.replace(0, 1)
    sp = (mo_sale/denom*100).values; pp = (mo_prop/denom*100).values
    x2 = np.arange(len(s.MONTHS)); w2 = 0.5
    ax2.bar(x2, pp, w2, label='プロパー品', color=C_TEAL, alpha=0.85, zorder=3)
    ax2.bar(x2, sp, w2, bottom=pp, label='在庫区分Dアイテム', color=C_RED, alpha=0.75, zorder=3)
    for i,(p_,sp_) in enumerate(zip(pp, sp)):
        _label_stacked(ax2, i, p_, sp_, fs_prop=12)
    ax2.set_xticks(x2); ax2.set_xticklabels(s.MONTHS, fontsize=13)
    ax2.set_ylabel('構成比（%）', fontsize=12); ax2.set_ylim(0, 125); ax2.legend(fontsize=11)
    ax2.set_title(f'{s.year26}年 月別 在庫区分構成', fontsize=12)
    ax2.grid(axis='y', alpha=0.3, zorder=1); ax2.spines[['top','right']].set_visible(False)
    fig.tight_layout(pad=1.5); return _to_buf(fig)

def ch_sale_products(sale_prods26, sale_qty_26):
    names=list(sale_prods26.index); qtys=list(sale_prods26['数量'])
    shares=[q/sale_qty_26*100 if sale_qty_26>0 else 0 for q in qtys]
    fig, ax = plt.subplots(figsize=(9.0, 4.5), facecolor='white')
    y = np.arange(len(names))
    ax.barh(y, shares, color=C_RED, alpha=0.75, height=0.65, zorder=3)
    ax.set_yticks(y); ax.set_yticklabels(names, fontsize=12)
    for i,(v,q) in enumerate(zip(shares,qtys)):
        ax.text(v+0.3, i, f'{v:.1f}%', va='center', fontsize=12, color=C_DARK)
    ax.set_xlabel('在庫区分Dアイテム内 構成比（%）', fontsize=12)
    ax.set_xlim(0, max(shares)+10); ax.invert_yaxis()
    ax.grid(axis='x', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(); return _to_buf(fig)

def ch_site(d25, d26):
    s25 = d25.groupby('得意先').agg(数量=('売上数','sum'), 金額=('売上金額','sum'))
    s26 = d26.groupby('得意先').agg(数量=('売上数','sum'), 金額=('売上金額','sum'))
    all_sites = sorted(set(s25.index)|set(s26.index),
                       key=lambda si: s26.loc[si,'金額'] if si in s26.index else 0, reverse=True)
    tot25=d25['売上金額'].sum(); tot26=d26['売上金額'].sum()
    sh25=[s25.loc[si,'金額']/tot25*100 if si in s25.index and tot25>0 else 0 for si in all_sites]
    sh26=[s26.loc[si,'金額']/tot26*100 if si in s26.index and tot26>0 else 0 for si in all_sites]
    r25=[s25.loc[si,'金額'] if si in s25.index else 0 for si in all_sites]
    r26=[s26.loc[si,'金額'] if si in s26.index else 0 for si in all_sites]
    growth=[((b/a-1)*100 if a>0 else None) for a,b in zip(r25,r26)]
    fig, (ax1, ax2) = plt.subplots(1, 2,
                                    figsize=(13, max(4.5, len(all_sites)*1.1+1.5)), facecolor='white')
    y=np.arange(len(all_sites)); w=0.35
    ax1.barh(y-w/2, sh25, w, label='前年', color=C_GRAY, alpha=0.65, zorder=3)
    ax1.barh(y+w/2, sh26, w, label='今年', color=C_TEAL, alpha=0.85, zorder=3)
    for i,(v25_,v26_) in enumerate(zip(sh25,sh26)):
        if v25_>0: ax1.text(v25_+0.3, i-w/2, f'{v25_:.1f}%', va='center', fontsize=10, color=C_DARK)
        if v26_>0: ax1.text(v26_+0.3, i+w/2, f'{v26_:.1f}%', va='center', fontsize=10,
                            color=C_TEAL, fontweight='bold')
    ax1.set_yticks(y); ax1.set_yticklabels(all_sites, fontsize=12)
    ax1.set_xlabel('売上高 構成比（%）', fontsize=11); ax1.legend(fontsize=11); ax1.invert_yaxis()
    ax1.set_xlim(0, max(sh25+sh26)+15)
    ax1.grid(axis='x', alpha=0.3, zorder=1); ax1.spines[['top','right']].set_visible(False)
    ax1.set_title('サイト別 売上高 構成比', fontsize=12, pad=8)
    g_vals=[g if g is not None else 0 for g in growth]
    ax2.barh(y, g_vals, color=[C_GREEN if g>=0 else C_RED for g in g_vals], alpha=0.85, height=0.5, zorder=3)
    for i,(g,g0) in enumerate(zip(g_vals,growth)):
        label=f'{g:+.1f}%' if g0 is not None else '新規'
        ax2.text(g+1.5 if g>=0 else g-1.5, i, label, va='center',
                 ha='left' if g>=0 else 'right', fontsize=11, fontweight='bold', color=C_DARK)
    ax2.axvline(0, color=C_DARK, lw=0.8, zorder=4)
    ax2.set_yticks(y); ax2.set_yticklabels(all_sites, fontsize=12)
    ax2.set_xlabel('前年比（%）', fontsize=11); ax2.invert_yaxis()
    ax2.grid(axis='x', alpha=0.3, zorder=1); ax2.spines[['top','right']].set_visible(False)
    ax2.set_title('サイト別 前年比成長率（売上高）', fontsize=12, pad=8)
    fig.tight_layout(); return _to_buf(fig)

def ch_march_index(months, unit_25_mo, unit_26_mo, rev_25_mo, rev_26_mo):
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 5.0), facecolor='white')
    for ax, d25m, d26m, title in [
        (ax1, unit_25_mo, unit_26_mo, f'販売点数（{months[0]}=100）'),
        (ax2, rev_25_mo,  rev_26_mo,  f'売上高（{months[0]}=100）'),
    ]:
        base=[v/d25m[0]*100 if d25m[0]>0 else 0 for v in d25m]
        curr=[v/d25m[0]*100 if d25m[0]>0 else 0 for v in d26m]
        x=np.arange(len(months)); w=0.35
        ax.bar(x-w/2, base, w, label='前年', color=C_GRAY, alpha=0.65, zorder=3)
        b2=ax.bar(x+w/2, curr, w, color=[C_RED if v<100 else C_GREEN for v in curr], alpha=0.85, zorder=3)
        for bar, v in zip(b2, curr):
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1.5,
                    f'{v:.0f}', ha='center', fontsize=12, fontweight='bold', color=C_DARK)
        ax.axhline(100, color=C_DARK, lw=0.8, ls='--', alpha=0.5, label='前年同水準')
        ax.set_xticks(x); ax.set_xticklabels(months, fontsize=13)
        ax.set_title(title, fontsize=12); ax.legend(fontsize=11)
        ax.set_ylim(0, max(curr)+35)
        ax.grid(axis='y', alpha=0.3, zorder=1); ax.spines[['top','right']].set_visible(False)
    fig.tight_layout(); return _to_buf(fig)

# ══════════════════════════════════════
# PPTX ユーティリティ
# ══════════════════════════════════════
def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    return prs

def _blank(prs):
    bl = next((l for l in prs.slide_layouts if l.name=='Blank'), prs.slide_layouts[-1])
    return prs.slides.add_slide(bl)

def _rect(sl, l, t, w, h, fill=None, line=None, lw=1):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    else:    sh.fill.background()
    if line: sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def _txt(sl, text, l, t, w, h, size=14, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = rgb(color); run.font.name = 'Yu Gothic'
    return txb

def _img(sl, buf, l, t, w, h):
    buf.seek(0); sl.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))

def _hdr(sl, title, sub=None, bar=C_TEAL):
    _rect(sl, 0, 0, 13.33, 1.05, fill=bar)
    _txt(sl, title, 0.4, 0.1, 12.0, 0.65, size=23, bold=True, color=C_WHITE)
    if sub: _txt(sl, sub, 0.4, 0.68, 12.0, 0.38, size=12, color='#DDEEEE')
    _rect(sl, 0, 7.32, 13.33, 0.18, fill=bar)

# ══════════════════════════════════════
# スライド関数
# ══════════════════════════════════════
def _s_cover(prs, s, note=''):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
    _rect(sl, 0, 0, 0.55, 7.5, fill=C_TEAL)
    _rect(sl, 0.55, 3.6, 12.78, 0.06, fill=C_GOLD)
    _txt(sl, f'{s.brand}  EC 販売実績分析レポート',
         1.0, 1.4, 11.8, 1.3, size=33, bold=True, color=C_WHITE)
    _txt(sl, f'{s.year25}年 {s.period}  vs  {s.year26}年 {s.period}  ―  前年同期比較{note}',
         1.0, 2.8, 11.8, 0.7, size=18, color=C_TEAL)
    _txt(sl, f'{s.year26}年　　CONFIDENTIAL  ─  社外秘',
         1.0, 4.1, 8.0, 0.55, size=13, color='#888888')

def _s_highlights(prs, s, tot_u, tot_r, rev_g, new_colors, label=''):
    sl = _blank(prs)
    _hdr(sl, f'ハイライト  ―  {s.year26}年 {s.period} 累計{label}')
    u_s = f'+{tot_u:.1f}%' if tot_u>=0 else f'{tot_u:.1f}%'
    r_s = f'+{tot_r:.1f}%' if tot_r>=0 else f'{tot_r:.1f}%'
    kpis=[
        (r_s, '売上高\n前年比成長率', C_TEAL, f'{s.MONTHS[-1]}が期間全体をリード'),
        (u_s, '販売点数\n前年比成長率', C_GREEN if tot_u>=0 else C_RED,
                                       f'{s.MONTHS[-1]}は{rev_g[-1]:+.1f}%（売上高）'),
        (f'+{new_colors}色', f'{s.year26}年\n新色追加数', C_GOLD, 'カラーバリエーション拡充'),
    ]
    for i,(val,lbl,col,note) in enumerate(kpis):
        l = 0.45+i*4.28
        _rect(sl, l, 1.35, 3.95, 3.9, fill=C_LGRAY, line=col, lw=2)
        _rect(sl, l, 1.35, 3.95, 0.28, fill=col)
        _txt(sl, val,  l, 1.85, 3.95, 1.3, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txt(sl, lbl,  l, 3.1,  3.95, 0.8, size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        _txt(sl, note, l, 3.9,  3.95, 0.9, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

def _s_monthly(prs, s, ug, rg, tot_u, tot_r, sub=''):
    sl = _blank(prs)
    _hdr(sl, '月別 前年比成長率', sub=sub or '販売点数・売上高の月次推移')
    _img(sl, ch_monthly(ug, rg, tot_u, tot_r, s.MONTHS), 0.7, 1.1, 11.9, 6.0)

def _s_category(prs, cat25, cat26, sub=''):
    sl = _blank(prs)
    _hdr(sl, '商品カテゴリ別 構成比', sub=sub or 'カテゴリ別の需要構造と前年からの変化')
    _img(sl, ch_category(cat25, cat26), 1.5, 1.2, 10.3, 5.0)

def _s_ranking(prs, top10, sub=''):
    sl = _blank(prs)
    _hdr(sl, '売れ筋商品ランキング  TOP10（今年）', sub=sub or '上位商品の構成比と前年との比較')
    _img(sl, ch_ranking(top10), 0.7, 1.1, 11.9, 6.1)

def _s_color(prs, s, top_colors, new_colors, sub=''):
    sl = _blank(prs)
    _hdr(sl, 'カラートレンド', sub=sub or f'新色 {new_colors}色を追加')
    _img(sl, ch_color(top_colors, new_colors), 0.7, 1.1, 11.9, 6.1)

def _s_size(prs, sz_labels, sz_2025, sz_2026, sub=''):
    sl = _blank(prs)
    _hdr(sl, 'サイズ構成', sub=sub or '2年連続でほぼ同一の安定分布')
    _img(sl, ch_size(sz_labels, sz_2025, sz_2026), 0.7, 1.1, 11.9, 5.9)
    _txt(sl, "※ サイズ略称（例：M5W7 = Men's 5 / Women's 7）",
         0.6, 7.06, 12.5, 0.35, size=11, color=C_GRAY)

def _s_summary_pos(prs, s, tot_u, tot_r, rev_g, unit_g, new_colors):
    sl = _blank(prs)
    _hdr(sl, 'まとめ  ―  展示会向けポジティブポイント', bar=C_GREEN)
    u_s = f'+{tot_u:.1f}%' if tot_u>=0 else f'{tot_u:.1f}%'
    r_s = f'+{tot_r:.1f}%' if tot_r>=0 else f'{tot_r:.1f}%'
    m1 = s.MONTHS[1] if len(s.MONTHS)>1 else s.MONTHS[0]
    pts=[
        (f'売上高 {r_s}  販売点数 {u_s}', f'{s.period}の累計で前年を上回るペースで推移。'),
        (f'{s.MONTHS[-1]} 売上高 {rev_g[-1]:+.1f}%', '期間後半に成長が加速。ブランド力を裏付ける数字。'),
        (f'{m1} 売上高 {rev_g[1 if len(rev_g)>1 else 0]:+.1f}%', '中盤以降は明確なプラス成長に転換。'),
        (f'新色 {new_colors}色 追加', '新鮮なカラーラインナップで顧客の購買意欲を維持・拡大。'),
        ('サイズ需要が安定', '2年連続でサイズ構成比がほぼ一致。在庫計画が立てやすい。'),
        ('カテゴリの多様化', '幅広いシーンでの需要を確認。複数カテゴリでの拡大傾向。'),
    ]
    for i,(ttl,body) in enumerate(pts):
        row=i//2; col=i%2
        l=0.4+col*6.5; t=1.3+row*1.9
        _rect(sl, l, t, 6.1, 1.75, fill=C_LGRAY, line=C_GREEN, lw=1)
        _txt(sl, ttl,  l+0.15, t+0.1,  5.8, 0.55, size=13, bold=True, color=C_DARK)
        _txt(sl, body, l+0.15, t+0.62, 5.8, 1.0,  size=11.5, color='#333333')

def _s_portfolio(prs, s):
    sl = _blank(prs)
    _hdr(sl, '【社内参考】商品ラインナップ拡充（プロパー品）',
         sub='モデル数・SKU数・カラー数の前年比較', bar=C_PURPLE)
    _img(sl, ch_portfolio(s.PORT, s.year25, s.year26), 1.5, 1.2, 10.3, 5.1)
    p25c=s.PORT[str(s.year25)]; p26c=s.PORT[str(s.year26)]
    _rect(sl, 0.4, 6.45, 12.5, 0.75, fill='#E8F9FA', line=C_TEAL, lw=1)
    _txt(sl, f'新カラー {s.NEW_COLORS_P}色  ／  SKU {p25c[1]}→{p26c[1]}（+{p26c[1]-p25c[1]}）  ／  カラー数 {p25c[2]}→{p26c[2]}（+{p26c[2]-p25c[2]}）',
         0.6, 6.55, 12.2, 0.6, size=13.5, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

def _s_site(prs, s):
    sl = _blank(prs)
    _hdr(sl, '【内部】サイト別分析', sub='得意先別 売上高 構成比・前年比成長率（全データ）', bar=C_PURPLE)
    _img(sl, ch_site(s.d25, s.d26), 0.2, 1.1, 12.9, 6.1)

def _s_divider(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill='#2C3E50')
    _rect(sl, 0, 0, 0.55, 7.5, fill=C_RED)
    _txt(sl, '内部確認事項', 1.0, 2.3, 11, 1.2, size=38, bold=True, color=C_WHITE)
    _txt(sl, '以下のスライドは社内共有用です。卸先には提示しないでください。',
         1.0, 3.8, 11, 0.7, size=16, color='#BBBBBB')

def _s_sale_analysis(prs, s):
    sl = _blank(prs)
    _hdr(sl, '【内部】在庫区分Dアイテム分析',
         sub=f'{s.year26}年に発生。プロパー品との比率と内訳を確認', bar=C_PURPLE)
    _img(sl, ch_sale_ratio(s), 0.5, 1.15, 12.3, 5.3)
    _rect(sl, 0.4, 6.55, 12.5, 0.7, fill='#FDEDEC', line=C_RED, lw=1)
    _txt(sl, f'{s.year25}年にはゼロ → {s.year26}年は全体の{s.sale_ratio_26:.1f}%（{s.sale_qty_26:.0f}点）発生。プロパー品ベースの成長率と全体数値の乖離に注意。',
         0.6, 6.63, 12.3, 0.58, size=12, color='#7B241C')

def _s_sale_products(prs, s):
    sl = _blank(prs)
    _hdr(sl, f'【内部】在庫区分Dアイテム 商品別内訳（{s.year26}年）', bar=C_PURPLE)
    _img(sl, ch_sale_products(s.sale_prods26, s.sale_qty_26), 0.7, 1.2, 11.9, 5.0)
    avg_s = s.d26_sale['売上金額'].sum()/s.sale_qty_26 if s.sale_qty_26>0 else 0
    avg_p = s.p26['売上金額'].sum()/s.p26['売上数'].sum() if s.p26['売上数'].sum()>0 else 0
    _rect(sl, 0.4, 6.4, 12.5, 0.85, fill='#FDEDEC', line=C_RED, lw=1)
    _txt(sl, f'区分D平均単価: {avg_s:,.0f}円  vs  プロパー品平均単価: {avg_p:,.0f}円。どのモデルが対象か確認し、仕入れ・展開方針に活用。',
         0.6, 6.48, 12.3, 0.73, size=12, color='#7B241C')

def _s_march_internal(prs, s):
    sl = _blank(prs)
    _hdr(sl, '【内部】月別動向の詳細（全データ・指数）',
         sub=f'{s.MONTHS[0]}を100とした指数で月次の動向を確認', bar=C_PURPLE)
    _img(sl, ch_march_index(s.MONTHS, s.UNIT_25_MO, s.UNIT_26_MO, s.REV_25_MO, s.REV_26_MO),
         0.5, 1.15, 12.3, 5.3)
    note = f'{s.MONTHS[0]}：数量{s.UNIT_G_A[0]:+.1f}%・売上高{s.REV_G_A[0]:+.1f}%。'
    if len(s.MONTHS) > 1: note += f'{s.MONTHS[1]}{s.REV_G_A[1]:+.1f}%'
    if len(s.MONTHS) > 2: note += f' → {s.MONTHS[2]}{s.REV_G_A[2]:+.1f}%と急回復・加速。'
    _rect(sl, 0.4, 6.55, 12.5, 0.7, fill='#EAF4FB', line=C_NAVY, lw=1)
    _txt(sl, note, 0.6, 6.63, 12.3, 0.58, size=12, color='#1A376C')

def _s_internal_summary(prs, s):
    sl = _blank(prs)
    _hdr(sl, '【内部】課題整理', bar=C_PURPLE)
    site_note = ''
    if s.worst_yoy < 0:
        site_note += f'{s.worst_site} が数量前年比{s.worst_yoy:+.1f}%と減少。'
    if s.site_gone: site_note += f'停止サイト：{s.site_gone}。'
    if s.site_new:  site_note += f'新規稼働：{s.site_new}。'
    site_note += 'チャネル間の需要移動か実需の減少かを確認し、各サイトの品揃え・販促戦略を検討。'
    p25c=s.PORT[str(s.year25)]; p26c=s.PORT[str(s.year26)]
    items=[
        (C_RED,  f'{s.MONTHS[0]}動向の要因確認',
         f'{s.MONTHS[0]}：数量{s.UNIT_G_A[0]:+.1f}%・売上高{s.REV_G_A[0]:+.1f}%。在庫切れ・価格変更等の要因を精査。'),
        (C_RED,  f'在庫区分Dアイテムの発生（全体の{s.sale_ratio_26:.1f}%）',
         f'{s.year25}年にはなかった在庫区分Dアイテムが{s.year26}年に{s.sale_qty_26:.0f}点発生。消化目的か継続するかの方針整理が必要。卸先への提示は不要。'),
        (C_NAVY, 'サイト別チャネル動向の確認', site_note),
        (C_GOLD, 'ラインナップ拡充（社内参考）',
         f'プロパー品でモデル数+{p26c[0]-p25c[0]}・SKU+{p26c[1]-p25c[1]}・新色{s.NEW_COLORS_P}色。卸先は全ライン非取扱のため展示会では訴求不要。'),
    ]
    for i,(col,ttl,body) in enumerate(items):
        t=1.3+i*1.42
        _rect(sl, 0.4, t, 12.5, 1.3, fill=C_LGRAY, line=col, lw=2)
        _rect(sl, 0.4, t, 0.28, 1.3, fill=col)
        _txt(sl, ttl,  0.8, t+0.1,  11.8, 0.5,  size=14, bold=True, color=C_DARK)
        _txt(sl, body, 0.8, t+0.58, 11.8, 0.68, size=12, color='#333333')

# ══════════════════════════════════════
# レポート生成（BytesIO を返す）
# ══════════════════════════════════════
def make_pos(s):
    """展示会用（全データ・区分表記なし）"""
    prs = _new_prs()
    _s_cover(prs, s)
    _s_highlights(prs, s, s.TOT_UNIT_A, s.TOT_REV_A, s.REV_G_A, s.NEW_COLORS_A)
    _s_monthly(prs, s, s.UNIT_G_A, s.REV_G_A, s.TOT_UNIT_A, s.TOT_REV_A)
    _s_category(prs, s.CAT25_A, s.CAT26_A)
    _s_ranking(prs, s.TOP10_A)
    _s_color(prs, s, s.TOP_COLORS_A, s.NEW_COLORS_A)
    _s_size(prs, s.SZ_LABELS_A, s.SZ_2025_A, s.SZ_2026_A)
    _s_summary_pos(prs, s, s.TOT_UNIT_A, s.TOT_REV_A, s.REV_G_A, s.UNIT_G_A, s.NEW_COLORS_A)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf

def make_int(s):
    """社内版（全データ＋在庫区分D分析＋サイト別）"""
    prs = _new_prs()
    _s_cover(prs, s)
    _s_highlights(prs, s, s.TOT_UNIT_A, s.TOT_REV_A, s.REV_G_A, s.NEW_COLORS_A)
    _s_monthly(prs, s, s.UNIT_G_A, s.REV_G_A, s.TOT_UNIT_A, s.TOT_REV_A)
    _s_category(prs, s.CAT25_A, s.CAT26_A)
    _s_ranking(prs, s.TOP10_A)
    _s_color(prs, s, s.TOP_COLORS_A, s.NEW_COLORS_A)
    _s_size(prs, s.SZ_LABELS_A, s.SZ_2025_A, s.SZ_2026_A)
    _s_summary_pos(prs, s, s.TOT_UNIT_A, s.TOT_REV_A, s.REV_G_A, s.UNIT_G_A, s.NEW_COLORS_A)
    _s_portfolio(prs, s)
    _s_site(prs, s)
    _s_divider(prs)
    _s_sale_analysis(prs, s)
    _s_sale_products(prs, s)
    _s_march_internal(prs, s)
    _s_internal_summary(prs, s)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf
